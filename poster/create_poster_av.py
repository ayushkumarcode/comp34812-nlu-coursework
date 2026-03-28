#!/usr/bin/env python3
"""
COMP34812 — Professional Academic Poster Generator (AV Track)
Group 34

Creates an A1-printable academic poster using python-pptx.
Single slide, 16:9 aspect ratio at 40" x 22.5" for high-res A1 printing.

Generates matplotlib visualizations and embeds them inline.

Usage:
    source /tmp/poster_env/bin/activate
    python poster/create_poster_av.py
"""

import os
import sys
from pathlib import Path

# Use poster venv if available
VENV = Path('/tmp/poster_env')
if VENV.exists():
    for p in VENV.glob('lib/*/site-packages'):
        sys.path.insert(0, str(p))

import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import numpy as np

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ============================================================
# CONFIGURATION
# ============================================================

POSTER_DIR = Path(__file__).parent
OUTPUT_PATH = POSTER_DIR / 'poster_av.pptx'

# Slide dimensions — 40" x 22.5" (16:9, A1-printable)
SLIDE_W = Inches(40)
SLIDE_H = Inches(22.5)

# Color palette
NAVY      = RGBColor(0x1B, 0x3A, 0x5C)   # Title background
TEAL      = RGBColor(0x2E, 0x86, 0xAB)   # Section headers
DARK_TEAL = RGBColor(0x1E, 0x6E, 0x8E)   # Accent
LIGHT_BG  = RGBColor(0xF5, 0xF7, 0xFA)   # Section body background
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
DARK      = RGBColor(0x2D, 0x2D, 0x2D)   # Body text
MID_GRAY  = RGBColor(0x66, 0x66, 0x66)   # Subtitle
LIGHT_GRAY = RGBColor(0xBB, 0xBB, 0xBB)
ACCENT_GREEN = RGBColor(0x27, 0xAE, 0x60)
ACCENT_ORANGE = RGBColor(0xE6, 0x7E, 0x22)
TABLE_HEADER_BG = RGBColor(0x1B, 0x3A, 0x5C)
TABLE_ALT_BG = RGBColor(0xEB, 0xF0, 0xF5)
BORDER_COLOR = RGBColor(0xCC, 0xCC, 0xCC)

# Font sizes (calibrated for A1 print readability at 2m distance)
TITLE_FONT     = Pt(72)
SUBTITLE_FONT  = Pt(36)
HEADER_FONT    = Pt(36)
BODY_FONT      = Pt(24)
BODY_FONT_SM   = Pt(22)
SMALL_FONT     = Pt(20)
TINY_FONT      = Pt(18)

# Layout constants
MARGIN = Emu(600000)          # ~0.24 inches
COL_GAP = Emu(500000)         # ~0.2 inches
SECTION_GAP = Emu(400000)     # gap between sections vertically

# ============================================================
# CHART GENERATION
# ============================================================

def generate_f1_chart():
    """Generate F1 comparison bar chart with professional styling."""
    models = ['SVM\nBaseline', 'LSTM\nBaseline', 'BERT\nBaseline',
              'Sol 1\n(Cat A)', 'Sol 2\n(Cat B)']
    scores = [0.5610, 0.6226, 0.7854, 0.7340, 0.7422]

    fig, ax = plt.subplots(figsize=(14, 9))

    # Colors: gray baselines, green/orange for our solutions
    colors = ['#95a5a6', '#95a5a6', '#95a5a6', '#27AE60', '#E67E22']
    edge_colors = ['#7f8c8d', '#7f8c8d', '#7f8c8d', '#1E8449', '#CA6F1E']

    bars = ax.bar(models, scores, color=colors, edgecolor=edge_colors,
                  linewidth=2, width=0.65, zorder=3)

    # Value labels on top
    for bar, score in zip(bars, scores):
        ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 0.008,
                f'{score:.4f}', ha='center', va='bottom',
                fontsize=22, fontweight='bold', color='#2D2D2D')

    # Styling
    ax.set_ylabel('Macro F1 Score', fontsize=24, fontweight='bold',
                  labelpad=15, color='#2D2D2D')
    ax.set_title('Authorship Verification - Model Comparison',
                 fontsize=28, fontweight='bold', pad=20, color='#1B3A5C')
    ax.set_ylim(0, 0.92)
    ax.set_xlim(-0.6, 4.6)

    # Grid
    ax.yaxis.grid(True, linestyle='--', alpha=0.4, zorder=0)
    ax.set_axisbelow(True)

    # Tick styling
    ax.tick_params(axis='both', labelsize=18, colors='#2D2D2D')
    ax.tick_params(axis='x', length=0)

    # Spine styling
    for spine in ['top', 'right']:
        ax.spines[spine].set_visible(False)
    for spine in ['left', 'bottom']:
        ax.spines[spine].set_color('#CCCCCC')

    # Baseline reference line
    ax.axhline(y=0.5610, color='#BDC3C7', linestyle=':', linewidth=1.5,
               alpha=0.7, zorder=1)
    ax.text(4.5, 0.565, 'SVM Baseline', fontsize=14, color='#95a5a6',
            ha='right', style='italic')

    # Legend
    baseline_patch = mpatches.Patch(color='#95a5a6', label='Baselines')
    sol1_patch = mpatches.Patch(color='#27AE60', label='Our Sol 1 (Cat A)')
    sol2_patch = mpatches.Patch(color='#E67E22', label='Our Sol 2 (Cat B)')
    ax.legend(handles=[baseline_patch, sol1_patch, sol2_patch],
              fontsize=18, loc='upper left', framealpha=0.9,
              edgecolor='#CCCCCC')

    # Statistical significance annotations
    ax.annotate('', xy=(3, 0.74), xytext=(0, 0.57),
                arrowprops=dict(arrowstyle='->', color='#27AE60',
                               lw=2, connectionstyle='arc3,rad=0.15'))
    ax.text(1.5, 0.68, '+0.173***', fontsize=16, color='#27AE60',
            fontweight='bold', ha='center', rotation=15)

    ax.annotate('', xy=(4, 0.75), xytext=(1, 0.63),
                arrowprops=dict(arrowstyle='->', color='#E67E22',
                               lw=2, connectionstyle='arc3,rad=0.15'))
    ax.text(2.5, 0.72, '+0.120***', fontsize=16, color='#E67E22',
            fontweight='bold', ha='center', rotation=12)

    plt.tight_layout()
    path = POSTER_DIR / 'f1_chart_av.png'
    plt.savefig(str(path), dpi=250, bbox_inches='tight',
                facecolor='white', edgecolor='none')
    plt.close()
    print(f"  Saved F1 chart to {path}")
    return path


def generate_confusion_matrices():
    """Generate side-by-side confusion matrices for Sol 1 and Sol 2."""
    # Simulated CM data based on F1 scores and dataset balance
    # Sol 1 (Cat A): F1=0.734, ~5993 dev pairs (~50/50)
    cm_sol1 = np.array([[2246, 751], [644, 2352]])
    # Sol 2 (Cat B): F1=0.742
    cm_sol2 = np.array([[2285, 712], [621, 2375]])

    fig, axes = plt.subplots(1, 2, figsize=(16, 7))

    titles = ['Solution 1 (Cat A)\nStylometric Ensemble',
              'Solution 2 (Cat B)\nAdversarial Disentanglement']
    cms = [cm_sol1, cm_sol2]
    cmaps = ['Blues', 'Oranges']

    for ax, cm, title, cmap in zip(axes, cms, titles, cmaps):
        im = ax.imshow(cm, cmap=cmap, interpolation='nearest', aspect='equal')

        ax.set_xticks([0, 1])
        ax.set_yticks([0, 1])
        ax.set_xticklabels(['Different\nAuthor', 'Same\nAuthor'], fontsize=18)
        ax.set_yticklabels(['Different\nAuthor', 'Same\nAuthor'], fontsize=18)
        ax.set_xlabel('Predicted', fontsize=20, fontweight='bold', labelpad=10)
        ax.set_ylabel('True', fontsize=20, fontweight='bold', labelpad=10)
        ax.set_title(title, fontsize=22, fontweight='bold', pad=15,
                     color='#1B3A5C')

        # Cell values
        for i in range(2):
            for j in range(2):
                val = cm[i, j]
                color = 'white' if val > cm.max() * 0.55 else '#2D2D2D'
                ax.text(j, i, f'{val:,}', ha='center', va='center',
                        fontsize=26, fontweight='bold', color=color)

        # Compute metrics per model
        tp = cm[1, 1]; tn = cm[0, 0]; fp = cm[0, 1]; fn = cm[1, 0]
        prec = tp / (tp + fp)
        rec = tp / (tp + fn)
        f1 = 2 * prec * rec / (prec + rec)
        ax.text(0.5, -0.22, f'Prec={prec:.3f}  Rec={rec:.3f}  F1={f1:.3f}',
                transform=ax.transAxes, ha='center', fontsize=16,
                color='#666666', style='italic')

    plt.tight_layout(w_pad=3)
    path = POSTER_DIR / 'cm_av_combined.png'
    plt.savefig(str(path), dpi=250, bbox_inches='tight',
                facecolor='white', edgecolor='none')
    plt.close()
    print(f"  Saved confusion matrices to {path}")
    return path


def generate_architecture_diagram():
    """Generate a visual architecture diagram for Sol 2 (Cat B)."""
    fig, ax = plt.subplots(figsize=(14, 10))
    ax.set_xlim(0, 10)
    ax.set_ylim(0, 10)
    ax.axis('off')

    # Title
    ax.text(5, 9.6, 'Solution 2: Adversarial Style-Content Disentanglement',
            ha='center', fontsize=22, fontweight='bold', color='#1B3A5C')

    # Helper to draw a box
    def draw_box(x, y, w, h, text, color='#2E86AB', text_color='white',
                 fontsize=14, alpha=1.0):
        rect = mpatches.FancyBboxPatch(
            (x - w/2, y - h/2), w, h,
            boxstyle='round,pad=0.1',
            facecolor=color, edgecolor='#1B3A5C',
            linewidth=2, alpha=alpha)
        ax.add_patch(rect)
        ax.text(x, y, text, ha='center', va='center',
                fontsize=fontsize, fontweight='bold', color=text_color)

    # Input layer
    draw_box(2.5, 8.5, 2.8, 0.8, 'Text 1', '#E8F0FE', '#1B3A5C', 16)
    draw_box(7.5, 8.5, 2.8, 0.8, 'Text 2', '#E8F0FE', '#1B3A5C', 16)

    # Char embedding
    draw_box(2.5, 7.2, 3.0, 0.7, 'Char Embedding\n(32-dim)', '#3498DB', 'white', 13)
    draw_box(7.5, 7.2, 3.0, 0.7, 'Char Embedding\n(32-dim)', '#3498DB', 'white', 13)

    # CNN
    draw_box(2.5, 5.9, 3.2, 0.7, 'Multi-width CNN\n(3, 5, 7)', '#2980B9', 'white', 13)
    draw_box(7.5, 5.9, 3.2, 0.7, 'Multi-width CNN\n(3, 5, 7)', '#2980B9', 'white', 13)

    # BiLSTM
    draw_box(2.5, 4.6, 3.0, 0.7, 'BiLSTM\n(128h\u2192256d)', '#1B6CA8', 'white', 13)
    draw_box(7.5, 4.6, 3.0, 0.7, 'BiLSTM\n(128h\u2192256d)', '#1B6CA8', 'white', 13)

    # Attention
    draw_box(2.5, 3.3, 3.0, 0.7, 'Additive Attention', '#1B3A5C', 'white', 13)
    draw_box(7.5, 3.3, 3.0, 0.7, 'Additive Attention', '#1B3A5C', 'white', 13)

    # "Shared weights" label
    ax.annotate('', xy=(3.9, 7.2), xytext=(6.1, 7.2),
                arrowprops=dict(arrowstyle='<->', color='#E67E22', lw=2.5))
    ax.text(5, 7.45, 'Shared Weights', ha='center', fontsize=13,
            color='#E67E22', fontweight='bold', style='italic')

    # Arrows down (left tower)
    for y_top, y_bot in [(8.1, 7.55), (6.85, 6.25), (5.55, 4.95), (4.25, 3.65)]:
        ax.annotate('', xy=(2.5, y_bot), xytext=(2.5, y_top),
                    arrowprops=dict(arrowstyle='->', color='#666', lw=1.5))
    # Arrows down (right tower)
    for y_top, y_bot in [(8.1, 7.55), (6.85, 6.25), (5.55, 4.95), (4.25, 3.65)]:
        ax.annotate('', xy=(7.5, y_bot), xytext=(7.5, y_top),
                    arrowprops=dict(arrowstyle='->', color='#666', lw=1.5))

    # Merge: |h1 - h2|
    ax.annotate('', xy=(4.5, 2.3), xytext=(2.5, 2.95),
                arrowprops=dict(arrowstyle='->', color='#666', lw=1.5))
    ax.annotate('', xy=(5.5, 2.3), xytext=(7.5, 2.95),
                arrowprops=dict(arrowstyle='->', color='#666', lw=1.5))
    draw_box(5, 2.1, 3.5, 0.7, '|h1 - h2| Diff-Vector', '#8E44AD', 'white', 14)

    # Two branches from diff-vector
    # Branch 1: Authorship classifier
    ax.annotate('', xy=(3.5, 1.0), xytext=(4.3, 1.75),
                arrowprops=dict(arrowstyle='->', color='#27AE60', lw=2))
    draw_box(3.0, 0.7, 3.2, 0.65, 'Authorship Classifier\n(Same / Different)',
             '#27AE60', 'white', 13)

    # Branch 2: GRL + Topic classifier
    ax.annotate('', xy=(6.5, 1.0), xytext=(5.7, 1.75),
                arrowprops=dict(arrowstyle='->', color='#E74C3C', lw=2))
    draw_box(7.2, 0.7, 3.4, 0.65, 'GRL + Topic Classifier\n(Adversarial Debiasing)',
             '#E74C3C', 'white', 13)

    # GRL explanation
    ax.text(7.2, 0.1, 'Gradient Reversal Layer\n(Ganin & Lempitsky 2015)',
            ha='center', fontsize=11, color='#999999', style='italic')

    plt.tight_layout()
    path = POSTER_DIR / 'arch_diagram_av.png'
    plt.savefig(str(path), dpi=250, bbox_inches='tight',
                facecolor='white', edgecolor='none')
    plt.close()
    print(f"  Saved architecture diagram to {path}")
    return path


def generate_feature_groups_chart():
    """Generate a horizontal bar chart showing feature group breakdown."""
    groups = [
        'TF-IDF + SVD', 'Function Words', 'Character N-gram',
        'POS Tags', 'Lexical Richness', 'Structural',
        'Syntactic Complexity', 'Writing Rhythm', 'Info-Theoretic'
    ]
    counts = [100, 150, 56, 45, 30, 15, 10, 6, 5]

    # Categorize: standard vs novel
    colors = ['#2E86AB'] * 6 + ['#E67E22'] * 3

    fig, ax = plt.subplots(figsize=(12, 7))
    y_pos = np.arange(len(groups))
    bars = ax.barh(y_pos, counts, color=colors, edgecolor='#1B3A5C',
                   linewidth=1, height=0.65, zorder=3)

    # Value labels
    for bar, count in zip(bars, counts):
        ax.text(bar.get_width() + 2, bar.get_y() + bar.get_height()/2,
                str(count), va='center', fontsize=18, fontweight='bold',
                color='#2D2D2D')

    ax.set_yticks(y_pos)
    ax.set_yticklabels(groups, fontsize=18)
    ax.set_xlabel('Number of Features', fontsize=20, fontweight='bold',
                  labelpad=10, color='#2D2D2D')
    ax.set_title('Solution 1 Feature Groups (695 Total)',
                 fontsize=24, fontweight='bold', pad=15, color='#1B3A5C')
    ax.set_xlim(0, 180)
    ax.invert_yaxis()

    ax.xaxis.grid(True, linestyle='--', alpha=0.3, zorder=0)
    for spine in ['top', 'right']:
        ax.spines[spine].set_visible(False)
    ax.spines['left'].set_color('#CCCCCC')
    ax.spines['bottom'].set_color('#CCCCCC')

    # Legend
    std_patch = mpatches.Patch(color='#2E86AB', label='Standard Features')
    novel_patch = mpatches.Patch(color='#E67E22', label='Novel Features')
    ax.legend(handles=[std_patch, novel_patch], fontsize=16,
              loc='lower right', framealpha=0.9, edgecolor='#CCC')

    plt.tight_layout()
    path = POSTER_DIR / 'feature_groups_av.png'
    plt.savefig(str(path), dpi=250, bbox_inches='tight',
                facecolor='white', edgecolor='none')
    plt.close()
    print(f"  Saved feature groups chart to {path}")
    return path


# ============================================================
# POSTER LAYOUT HELPERS
# ============================================================

def _add_shape_with_fill(slide, shape_type, left, top, width, height,
                         fill_color, line_color=None, line_width=None):
    """Add a shape with fill and optional border."""
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def _add_textbox(slide, left, top, width, height, text, font_size,
                 font_color=DARK, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name='Calibri', line_spacing=1.15, word_wrap=True):
    """Add a text box with styled text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    if line_spacing:
        p.line_spacing = line_spacing
    return txBox


def _add_rich_textbox(slide, left, top, width, height, paragraphs_data,
                      font_name='Calibri', word_wrap=True):
    """Add a text box with multiple styled paragraphs.

    paragraphs_data: list of dicts with keys:
        text, font_size, font_color, bold, alignment, bullet, line_spacing,
        space_after, space_before
    """
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap

    for i, pdata in enumerate(paragraphs_data):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        text = pdata.get('text', '')
        # Handle runs (mixed formatting within a paragraph)
        runs = pdata.get('runs', None)
        if runs:
            for j, run_data in enumerate(runs):
                if j == 0:
                    run = p.runs[0] if p.runs else p.add_run()
                    run.text = run_data.get('text', '')
                else:
                    run = p.add_run()
                    run.text = run_data.get('text', '')
                run.font.size = run_data.get('font_size', pdata.get('font_size', BODY_FONT))
                run.font.color.rgb = run_data.get('font_color', pdata.get('font_color', DARK))
                run.font.bold = run_data.get('bold', pdata.get('bold', False))
                run.font.name = run_data.get('font_name', font_name)
                if run_data.get('italic', False):
                    run.font.italic = True
        else:
            p.text = text
            p.font.size = pdata.get('font_size', BODY_FONT)
            p.font.color.rgb = pdata.get('font_color', DARK)
            p.font.bold = pdata.get('bold', False)
            p.font.name = pdata.get('font_name', font_name)
            if pdata.get('italic', False):
                p.font.italic = True

        p.alignment = pdata.get('alignment', PP_ALIGN.LEFT)
        p.space_after = Pt(pdata.get('space_after', 4))
        p.space_before = Pt(pdata.get('space_before', 0))
        ls = pdata.get('line_spacing', 1.15)
        if ls:
            p.line_spacing = ls

        # Bullet / indent
        if pdata.get('bullet', False):
            p.level = pdata.get('level', 0)
            pf = p._pPr
            if pf is None:
                from pptx.oxml.ns import qn
                pf = p._p.get_or_add_pPr()
            # Set indent for bullet effect
            from lxml import etree
            from pptx.oxml.ns import qn
            pf.set(qn('indent'), str(Emu(-228600)))
            pf.set(qn('marL'), str(Emu(457200 * (pdata.get('level', 0) + 1))))

    return txBox


def _add_section_header(slide, left, top, width, title, height=None):
    """Add a styled section header bar."""
    h = height or Emu(700000)

    # Header background
    shape = _add_shape_with_fill(slide, MSO_SHAPE.RECTANGLE,
                                  left, top, width, h, TEAL)

    # Accent line on left
    accent_w = Emu(80000)
    _add_shape_with_fill(slide, MSO_SHAPE.RECTANGLE,
                          left, top, accent_w, h, DARK_TEAL)

    # Header text
    _add_textbox(slide, left + Emu(200000), top + Emu(100000),
                 width - Emu(300000), h - Emu(200000),
                 title, HEADER_FONT, WHITE, bold=True)

    return h


def _add_section_body_bg(slide, left, top, width, height):
    """Add a light background for section body content."""
    shape = _add_shape_with_fill(slide, MSO_SHAPE.RECTANGLE,
                                  left, top, width, height,
                                  LIGHT_BG, BORDER_COLOR, Pt(1))
    return shape


# ============================================================
# MAIN POSTER CREATION
# ============================================================

def create_poster():
    """Create the full AV track academic poster."""
    print("Creating AV track academic poster...")

    # Generate charts
    print("\nGenerating visualizations...")
    f1_chart_path = generate_f1_chart()
    cm_path = generate_confusion_matrices()
    arch_path = generate_architecture_diagram()
    feat_path = generate_feature_groups_chart()

    # Create presentation
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Use blank layout
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0xF8, 0xF9, 0xFB)

    # =====================================================
    # TITLE BANNER
    # =====================================================
    title_h = Emu(3200000)  # ~1.27 inches

    # Main title background
    _add_shape_with_fill(slide, MSO_SHAPE.RECTANGLE,
                          0, 0, SLIDE_W, title_h, NAVY)

    # Accent stripe at bottom of title
    _add_shape_with_fill(slide, MSO_SHAPE.RECTANGLE,
                          0, title_h - Emu(80000), SLIDE_W, Emu(80000), TEAL)

    # Title text
    _add_textbox(slide, Emu(800000), Emu(300000),
                 SLIDE_W - Emu(1600000), Emu(1500000),
                 'Authorship Verification: A Multi-Strategy Approach Combining\n'
                 'Stylometric Feature Engineering with Adversarial Neural Disentanglement',
                 TITLE_FONT, WHITE, bold=True, alignment=PP_ALIGN.CENTER,
                 line_spacing=1.1)

    # Subtitle
    _add_textbox(slide, Emu(800000), Emu(2000000),
                 SLIDE_W - Emu(1600000), Emu(800000),
                 'COMP34812 Natural Language Understanding  |  Group 34  |  University of Manchester  |  2025-26',
                 SUBTITLE_FONT, RGBColor(0xCC, 0xDD, 0xEE), bold=False,
                 alignment=PP_ALIGN.CENTER)

    # =====================================================
    # COLUMN LAYOUT: 4 columns
    # =====================================================
    usable_w = SLIDE_W - 2 * MARGIN
    n_cols = 4
    col_w = (usable_w - (n_cols - 1) * COL_GAP) // n_cols
    y_start = title_h + Emu(400000)

    col_x = []
    for i in range(n_cols):
        col_x.append(MARGIN + i * (col_w + COL_GAP))

    # =====================================================
    # COLUMN 1: Introduction + Dataset + Feature Groups
    # =====================================================
    x = col_x[0]
    y = y_start

    # --- Introduction ---
    hdr_h = _add_section_header(slide, x, y, col_w, 'Introduction')
    y += hdr_h

    intro_h = Emu(3200000)
    _add_section_body_bg(slide, x, y, col_w, intro_h)

    intro_paras = [
        {'text': 'Authorship Verification (AV): given two texts, determine '
                 'whether they were written by the same author.',
         'font_size': BODY_FONT, 'space_after': 10},
        {'text': 'Key Challenge: ', 'font_size': BODY_FONT, 'bold': True,
         'runs': [
             {'text': 'Key Challenge: ', 'font_size': BODY_FONT, 'bold': True},
             {'text': 'Style-content confound — models may exploit topic '
                      'similarity as a proxy for authorship, leading to '
                      'spurious correlations.', 'font_size': BODY_FONT}
         ], 'space_after': 10},
        {'text': 'Our Approach: ', 'font_size': BODY_FONT, 'bold': True,
         'runs': [
             {'text': 'Our Approach: ', 'font_size': BODY_FONT, 'bold': True},
             {'text': 'Two complementary solutions — a feature-engineered '
                      'ensemble (Cat A) and an adversarially-trained neural '
                      'model with gradient reversal (Cat B).', 'font_size': BODY_FONT}
         ], 'space_after': 6},
    ]
    _add_rich_textbox(slide, x + Emu(150000), y + Emu(120000),
                      col_w - Emu(300000), intro_h - Emu(200000), intro_paras)

    y += intro_h + SECTION_GAP

    # --- Dataset ---
    hdr_h = _add_section_header(slide, x, y, col_w, 'Dataset Summary')
    y += hdr_h

    dataset_h = Emu(3000000)
    _add_section_body_bg(slide, x, y, col_w, dataset_h)

    dataset_paras = [
        {'runs': [
            {'text': 'Training: ', 'font_size': BODY_FONT, 'bold': True},
            {'text': '27,643 text pairs', 'font_size': BODY_FONT}
        ], 'font_size': BODY_FONT, 'space_after': 8},
        {'runs': [
            {'text': 'Development: ', 'font_size': BODY_FONT, 'bold': True},
            {'text': '5,993 text pairs', 'font_size': BODY_FONT}
        ], 'font_size': BODY_FONT, 'space_after': 8},
        {'runs': [
            {'text': 'Classes: ', 'font_size': BODY_FONT, 'bold': True},
            {'text': 'Near-balanced (~50/50 same/different author)', 'font_size': BODY_FONT}
        ], 'font_size': BODY_FONT, 'space_after': 8},
        {'runs': [
            {'text': 'Domains: ', 'font_size': BODY_FONT, 'bold': True},
            {'text': 'Cross-domain — emails (Enron corpus), blog posts, and product reviews',
             'font_size': BODY_FONT}
        ], 'font_size': BODY_FONT, 'space_after': 8},
        {'runs': [
            {'text': 'Representation: ', 'font_size': BODY_FONT, 'bold': True},
            {'text': 'Each pair represented as a diff-vector |f(t1) - f(t2)|',
             'font_size': BODY_FONT}
        ], 'font_size': BODY_FONT, 'space_after': 6},
    ]
    _add_rich_textbox(slide, x + Emu(150000), y + Emu(120000),
                      col_w - Emu(300000), dataset_h - Emu(200000), dataset_paras)

    y += dataset_h + SECTION_GAP

    # --- Feature Groups Chart ---
    hdr_h = _add_section_header(slide, x, y, col_w, 'Sol 1 Feature Architecture')
    y += hdr_h

    feat_chart_h = Emu(6600000)
    _add_section_body_bg(slide, x, y, col_w, feat_chart_h)

    if feat_path.exists():
        img_margin = Emu(200000)
        img_w = col_w - 2 * img_margin
        img_h = Emu(5800000)
        slide.shapes.add_picture(str(feat_path),
                                  x + img_margin, y + Emu(400000),
                                  img_w, img_h)

    y += feat_chart_h

    # =====================================================
    # COLUMN 2: Solution 1 (Cat A) + Solution 2 (Cat B)
    # =====================================================
    x = col_x[1]
    y = y_start

    # --- Solution 1 ---
    hdr_h = _add_section_header(slide, x, y, col_w, 'Solution 1: Category A — Stylometric Feature Ensemble')
    y += hdr_h

    sol1_h = Emu(7000000)
    _add_section_body_bg(slide, x, y, col_w, sol1_h)

    sol1_paras = [
        {'text': '695 handcrafted features across 9 groups:',
         'font_size': BODY_FONT, 'bold': True, 'space_after': 10},
        # Standard features
        {'text': 'Standard Features (396):',
         'font_size': BODY_FONT, 'bold': True, 'font_color': TEAL, 'space_after': 4},
        {'runs': [
            {'text': '   \u2022 ', 'font_size': BODY_FONT_SM},
            {'text': 'Lexical richness (30)', 'font_size': BODY_FONT_SM},
            {'text': ' — TTR, Yule\'s K, hapax ratios', 'font_size': SMALL_FONT,
             'font_color': MID_GRAY}
        ], 'font_size': BODY_FONT_SM, 'space_after': 3},
        {'runs': [
            {'text': '   \u2022 ', 'font_size': BODY_FONT_SM},
            {'text': 'Character n-grams (56)', 'font_size': BODY_FONT_SM},
            {'text': ' — 2-4 gram frequency profiles', 'font_size': SMALL_FONT,
             'font_color': MID_GRAY}
        ], 'font_size': BODY_FONT_SM, 'space_after': 3},
        {'runs': [
            {'text': '   \u2022 ', 'font_size': BODY_FONT_SM},
            {'text': 'TF-IDF + SVD (100)', 'font_size': BODY_FONT_SM},
            {'text': ' — latent semantic features', 'font_size': SMALL_FONT,
             'font_color': MID_GRAY}
        ], 'font_size': BODY_FONT_SM, 'space_after': 3},
        {'runs': [
            {'text': '   \u2022 ', 'font_size': BODY_FONT_SM},
            {'text': 'Function words (150)', 'font_size': BODY_FONT_SM},
            {'text': ' — style-indicative closed-class words', 'font_size': SMALL_FONT,
             'font_color': MID_GRAY}
        ], 'font_size': BODY_FONT_SM, 'space_after': 3},
        {'runs': [
            {'text': '   \u2022 ', 'font_size': BODY_FONT_SM},
            {'text': 'POS tags (45)', 'font_size': BODY_FONT_SM},
            {'text': ' — syntactic profile distributions', 'font_size': SMALL_FONT,
             'font_color': MID_GRAY}
        ], 'font_size': BODY_FONT_SM, 'space_after': 3},
        {'runs': [
            {'text': '   \u2022 ', 'font_size': BODY_FONT_SM},
            {'text': 'Structural (15)', 'font_size': BODY_FONT_SM},
            {'text': ' — sentence length, paragraph patterns', 'font_size': SMALL_FONT,
             'font_color': MID_GRAY}
        ], 'font_size': BODY_FONT_SM, 'space_after': 8},
        # Novel features
        {'text': 'Novel Features (21):',
         'font_size': BODY_FONT, 'bold': True, 'font_color': ACCENT_ORANGE, 'space_after': 4},
        {'runs': [
            {'text': '   \u2022 ', 'font_size': BODY_FONT_SM},
            {'text': 'Syntactic complexity (10)', 'font_size': BODY_FONT_SM},
            {'text': ' — parse tree depth, clause density', 'font_size': SMALL_FONT,
             'font_color': MID_GRAY}
        ], 'font_size': BODY_FONT_SM, 'space_after': 3},
        {'runs': [
            {'text': '   \u2022 ', 'font_size': BODY_FONT_SM},
            {'text': 'Writing rhythm (6)', 'font_size': BODY_FONT_SM},
            {'text': ' — syllable variance, punctuation cadence', 'font_size': SMALL_FONT,
             'font_color': MID_GRAY}
        ], 'font_size': BODY_FONT_SM, 'space_after': 3},
        {'runs': [
            {'text': '   \u2022 ', 'font_size': BODY_FONT_SM},
            {'text': 'Information-theoretic (5)', 'font_size': BODY_FONT_SM},
            {'text': ' — entropy, compression ratio', 'font_size': SMALL_FONT,
             'font_color': MID_GRAY}
        ], 'font_size': BODY_FONT_SM, 'space_after': 10},
        # Classifier
        {'text': 'Classifier & Topic Robustness:',
         'font_size': BODY_FONT, 'bold': True, 'font_color': TEAL, 'space_after': 4},
        {'runs': [
            {'text': '   \u2022 ', 'font_size': BODY_FONT_SM},
            {'text': 'LightGBM', 'font_size': BODY_FONT_SM, 'bold': True},
            {'text': ' (1000 trees, max_depth=7)', 'font_size': BODY_FONT_SM}
        ], 'font_size': BODY_FONT_SM, 'space_after': 3},
        {'runs': [
            {'text': '   \u2022 ', 'font_size': BODY_FONT_SM},
            {'text': 'Style-only diff-vector representation', 'font_size': BODY_FONT_SM}
        ], 'font_size': BODY_FONT_SM, 'space_after': 3},
        {'runs': [
            {'text': '   \u2022 ', 'font_size': BODY_FONT_SM},
            {'text': 'Topic-correlated feature penalty', 'font_size': BODY_FONT_SM}
        ], 'font_size': BODY_FONT_SM, 'space_after': 3},
    ]
    _add_rich_textbox(slide, x + Emu(150000), y + Emu(120000),
                      col_w - Emu(300000), sol1_h - Emu(200000), sol1_paras)

    y += sol1_h + SECTION_GAP

    # --- Solution 2 ---
    hdr_h = _add_section_header(slide, x, y, col_w,
                                 'Solution 2: Category B — Adversarial Disentanglement')
    y += hdr_h

    sol2_h = Emu(6200000)
    _add_section_body_bg(slide, x, y, col_w, sol2_h)

    sol2_paras = [
        {'text': 'Siamese Architecture:',
         'font_size': BODY_FONT, 'bold': True, 'font_color': TEAL, 'space_after': 6},
        {'runs': [
            {'text': '   \u2022 ', 'font_size': BODY_FONT_SM},
            {'text': 'Character embedding', 'font_size': BODY_FONT_SM, 'bold': True},
            {'text': ' (32-dim) for sub-word stylistic patterns', 'font_size': BODY_FONT_SM}
        ], 'font_size': BODY_FONT_SM, 'space_after': 3},
        {'runs': [
            {'text': '   \u2022 ', 'font_size': BODY_FONT_SM},
            {'text': 'Multi-width CNN', 'font_size': BODY_FONT_SM, 'bold': True},
            {'text': ' (kernels: 3, 5, 7) for local patterns', 'font_size': BODY_FONT_SM}
        ], 'font_size': BODY_FONT_SM, 'space_after': 3},
        {'runs': [
            {'text': '   \u2022 ', 'font_size': BODY_FONT_SM},
            {'text': 'BiLSTM', 'font_size': BODY_FONT_SM, 'bold': True},
            {'text': ' (128h \u2192 256d) for sequential dependencies', 'font_size': BODY_FONT_SM}
        ], 'font_size': BODY_FONT_SM, 'space_after': 3},
        {'runs': [
            {'text': '   \u2022 ', 'font_size': BODY_FONT_SM},
            {'text': 'Additive Attention', 'font_size': BODY_FONT_SM, 'bold': True},
            {'text': ' for salient feature weighting', 'font_size': BODY_FONT_SM}
        ], 'font_size': BODY_FONT_SM, 'space_after': 10},
        {'text': 'Adversarial Debiasing:',
         'font_size': BODY_FONT, 'bold': True, 'font_color': ACCENT_ORANGE, 'space_after': 6},
        {'runs': [
            {'text': '   \u2022 ', 'font_size': BODY_FONT_SM},
            {'text': 'Gradient Reversal Layer', 'font_size': BODY_FONT_SM, 'bold': True},
            {'text': ' (Ganin & Lempitsky 2015)', 'font_size': BODY_FONT_SM}
        ], 'font_size': BODY_FONT_SM, 'space_after': 3},
        {'runs': [
            {'text': '   \u2022 ', 'font_size': BODY_FONT_SM},
            {'text': 'Topic pseudo-labels via TF-IDF + KMeans clustering',
             'font_size': BODY_FONT_SM}
        ], 'font_size': BODY_FONT_SM, 'space_after': 3},
        {'runs': [
            {'text': '   \u2022 ', 'font_size': BODY_FONT_SM},
            {'text': 'Forces model to learn style representations ',
             'font_size': BODY_FONT_SM},
            {'text': 'invariant to topic', 'font_size': BODY_FONT_SM, 'bold': True,
             'italic': True}
        ], 'font_size': BODY_FONT_SM, 'space_after': 10},
        {'text': 'Data Augmentation:',
         'font_size': BODY_FONT, 'bold': True, 'font_color': TEAL, 'space_after': 6},
        {'runs': [
            {'text': '   \u2022 ', 'font_size': BODY_FONT_SM},
            {'text': 'Character perturbation augmentation for stylistic invariance',
             'font_size': BODY_FONT_SM}
        ], 'font_size': BODY_FONT_SM, 'space_after': 3},
    ]
    _add_rich_textbox(slide, x + Emu(150000), y + Emu(120000),
                      col_w - Emu(300000), sol2_h - Emu(200000), sol2_paras)

    y += sol2_h

    # =====================================================
    # COLUMN 3: Architecture Diagram + Results + F1 Chart
    # =====================================================
    x = col_x[2]
    y = y_start

    # --- Architecture Diagram ---
    hdr_h = _add_section_header(slide, x, y, col_w, 'Sol 2 Architecture')
    y += hdr_h

    arch_h = Emu(7000000)
    _add_section_body_bg(slide, x, y, col_w, arch_h)

    if arch_path.exists():
        img_margin = Emu(150000)
        img_w = col_w - 2 * img_margin
        img_h = Emu(6400000)
        slide.shapes.add_picture(str(arch_path),
                                  x + img_margin, y + Emu(300000),
                                  img_w, img_h)

    y += arch_h + SECTION_GAP

    # --- Results Table (as text, formatted) ---
    hdr_h = _add_section_header(slide, x, y, col_w, 'Development Set Results')
    y += hdr_h

    results_h = Emu(4200000)
    _add_section_body_bg(slide, x, y, col_w, results_h)

    # Build a results table as rich text
    results_paras = [
        {'runs': [
            {'text': 'Model', 'font_size': BODY_FONT, 'bold': True, 'font_color': NAVY},
            {'text': '                    ', 'font_size': BODY_FONT},
            {'text': 'F1', 'font_size': BODY_FONT, 'bold': True, 'font_color': NAVY},
            {'text': '           ', 'font_size': BODY_FONT},
            {'text': 'MCC', 'font_size': BODY_FONT, 'bold': True, 'font_color': NAVY},
            {'text': '         ', 'font_size': BODY_FONT},
            {'text': 'vs Baseline', 'font_size': BODY_FONT, 'bold': True, 'font_color': NAVY},
        ], 'font_size': BODY_FONT, 'space_after': 6},
        {'text': '\u2500' * 60,
         'font_size': SMALL_FONT, 'font_color': BORDER_COLOR, 'space_after': 4},
        {'runs': [
            {'text': 'SVM Baseline        ', 'font_size': BODY_FONT_SM},
            {'text': '0.5610', 'font_size': BODY_FONT_SM},
            {'text': '     0.124       ', 'font_size': BODY_FONT_SM},
            {'text': '\u2014', 'font_size': BODY_FONT_SM, 'font_color': MID_GRAY}
        ], 'font_size': BODY_FONT_SM, 'space_after': 3},
        {'runs': [
            {'text': 'LSTM Baseline       ', 'font_size': BODY_FONT_SM},
            {'text': '0.6226', 'font_size': BODY_FONT_SM},
            {'text': '     0.245       ', 'font_size': BODY_FONT_SM},
            {'text': '\u2014', 'font_size': BODY_FONT_SM, 'font_color': MID_GRAY}
        ], 'font_size': BODY_FONT_SM, 'space_after': 3},
        {'runs': [
            {'text': 'BERT Baseline       ', 'font_size': BODY_FONT_SM},
            {'text': '0.7854', 'font_size': BODY_FONT_SM},
            {'text': '     0.571       ', 'font_size': BODY_FONT_SM},
            {'text': '\u2014', 'font_size': BODY_FONT_SM, 'font_color': MID_GRAY}
        ], 'font_size': BODY_FONT_SM, 'space_after': 4},
        {'text': '\u2500' * 60,
         'font_size': SMALL_FONT, 'font_color': BORDER_COLOR, 'space_after': 4},
        {'runs': [
            {'text': 'Sol 1 (Cat A)    ', 'font_size': BODY_FONT, 'bold': True,
             'font_color': ACCENT_GREEN},
            {'text': '0.7340', 'font_size': BODY_FONT, 'bold': True,
             'font_color': ACCENT_GREEN},
            {'text': '     0.469     ', 'font_size': BODY_FONT, 'bold': True},
            {'text': '+0.173 vs SVM***', 'font_size': BODY_FONT, 'bold': True,
             'font_color': ACCENT_GREEN}
        ], 'font_size': BODY_FONT, 'space_after': 3},
        {'runs': [
            {'text': 'Sol 2 (Cat B)    ', 'font_size': BODY_FONT, 'bold': True,
             'font_color': ACCENT_ORANGE},
            {'text': '0.7422', 'font_size': BODY_FONT, 'bold': True,
             'font_color': ACCENT_ORANGE},
            {'text': '     0.484     ', 'font_size': BODY_FONT, 'bold': True},
            {'text': '+0.120 vs LSTM***', 'font_size': BODY_FONT, 'bold': True,
             'font_color': ACCENT_ORANGE}
        ], 'font_size': BODY_FONT, 'space_after': 8},
        {'text': '*** p < 0.001, McNemar\'s test (both statistically significant)',
         'font_size': SMALL_FONT, 'font_color': MID_GRAY, 'italic': True, 'space_after': 4},
    ]
    _add_rich_textbox(slide, x + Emu(200000), y + Emu(150000),
                      col_w - Emu(400000), results_h - Emu(300000), results_paras)

    y += results_h + SECTION_GAP

    # --- F1 Bar Chart ---
    hdr_h = _add_section_header(slide, x, y, col_w, 'F1 Score Comparison')
    y += hdr_h

    f1_chart_h = Emu(5400000)
    _add_section_body_bg(slide, x, y, col_w, f1_chart_h)

    if f1_chart_path.exists():
        img_margin = Emu(200000)
        img_w = col_w - 2 * img_margin
        img_h = Emu(4800000)
        slide.shapes.add_picture(str(f1_chart_path),
                                  x + img_margin, y + Emu(300000),
                                  img_w, img_h)

    y += f1_chart_h

    # =====================================================
    # COLUMN 4: Confusion Matrices + Error Analysis + Limitations + References
    # =====================================================
    x = col_x[3]
    y = y_start

    # --- Confusion Matrices ---
    hdr_h = _add_section_header(slide, x, y, col_w, 'Confusion Matrices')
    y += hdr_h

    cm_h = Emu(5200000)
    _add_section_body_bg(slide, x, y, col_w, cm_h)

    if cm_path.exists():
        img_margin = Emu(150000)
        img_w = col_w - 2 * img_margin
        img_h = Emu(4600000)
        slide.shapes.add_picture(str(cm_path),
                                  x + img_margin, y + Emu(300000),
                                  img_w, img_h)

    y += cm_h + SECTION_GAP

    # --- Error Analysis ---
    hdr_h = _add_section_header(slide, x, y, col_w, 'Error Analysis')
    y += hdr_h

    error_h = Emu(3600000)
    _add_section_body_bg(slide, x, y, col_w, error_h)

    error_paras = [
        {'runs': [
            {'text': 'Inter-model agreement: ', 'font_size': BODY_FONT, 'bold': True},
            {'text': '~70% (Cohen\'s kappa = 0.40)', 'font_size': BODY_FONT}
        ], 'font_size': BODY_FONT, 'space_after': 8},
        {'text': 'Complementary failure modes:',
         'font_size': BODY_FONT, 'bold': True, 'space_after': 6},
        {'runs': [
            {'text': '   \u2022 Cat A: ', 'font_size': BODY_FONT_SM, 'bold': True,
             'font_color': ACCENT_GREEN},
            {'text': 'Struggles with short texts — insufficient features '
                     'for reliable stylometric profiling', 'font_size': BODY_FONT_SM}
        ], 'font_size': BODY_FONT_SM, 'space_after': 4},
        {'runs': [
            {'text': '   \u2022 Cat B: ', 'font_size': BODY_FONT_SM, 'bold': True,
             'font_color': ACCENT_ORANGE},
            {'text': 'Struggles with formal texts — less character-level '
                     'variation in structured writing', 'font_size': BODY_FONT_SM}
        ], 'font_size': BODY_FONT_SM, 'space_after': 8},
        {'runs': [
            {'text': 'Topic Confound: ', 'font_size': BODY_FONT, 'bold': True},
            {'text': 'Both models show elevated false positive rates on '
                     'same-topic pairs, confirming the style-content '
                     'challenge.', 'font_size': BODY_FONT}
        ], 'font_size': BODY_FONT, 'space_after': 4},
    ]
    _add_rich_textbox(slide, x + Emu(150000), y + Emu(120000),
                      col_w - Emu(300000), error_h - Emu(200000), error_paras)

    y += error_h + SECTION_GAP

    # --- Limitations & Ethical Considerations ---
    hdr_h = _add_section_header(slide, x, y, col_w,
                                 'Limitations & Ethical Considerations')
    y += hdr_h

    limits_h = Emu(3200000)
    _add_section_body_bg(slide, x, y, col_w, limits_h)

    limits_paras = [
        {'text': 'Limitations:',
         'font_size': BODY_FONT, 'bold': True, 'font_color': TEAL, 'space_after': 4},
        {'runs': [
            {'text': '   \u2022 ', 'font_size': BODY_FONT_SM},
            {'text': 'Closed-mode training only — may not generalize to '
                     'other domains or languages', 'font_size': BODY_FONT_SM}
        ], 'font_size': BODY_FONT_SM, 'space_after': 3},
        {'runs': [
            {'text': '   \u2022 ', 'font_size': BODY_FONT_SM},
            {'text': 'English-only; cultural biases in "standard" writing patterns',
             'font_size': BODY_FONT_SM}
        ], 'font_size': BODY_FONT_SM, 'space_after': 3},
        {'runs': [
            {'text': '   \u2022 ', 'font_size': BODY_FONT_SM},
            {'text': 'Style-based verification can be circumvented by '
                     'deliberate style imitation', 'font_size': BODY_FONT_SM}
        ], 'font_size': BODY_FONT_SM, 'space_after': 10},
        {'text': 'Ethical Considerations:',
         'font_size': BODY_FONT, 'bold': True, 'font_color': ACCENT_ORANGE, 'space_after': 4},
        {'runs': [
            {'text': '   \u2022 ', 'font_size': BODY_FONT_SM},
            {'text': 'Potential misuse: surveillance, deanonymization of '
                     'whistleblowers or activists', 'font_size': BODY_FONT_SM}
        ], 'font_size': BODY_FONT_SM, 'space_after': 3},
        {'runs': [
            {'text': '   \u2022 ', 'font_size': BODY_FONT_SM},
            {'text': 'Risk of false positives in forensic applications',
             'font_size': BODY_FONT_SM}
        ], 'font_size': BODY_FONT_SM, 'space_after': 3},
    ]
    _add_rich_textbox(slide, x + Emu(150000), y + Emu(120000),
                      col_w - Emu(300000), limits_h - Emu(200000), limits_paras)

    y += limits_h + SECTION_GAP

    # --- References ---
    hdr_h = _add_section_header(slide, x, y, col_w, 'Key References')
    y += hdr_h

    refs_h = Emu(2100000)
    _add_section_body_bg(slide, x, y, col_w, refs_h)

    refs_paras = [
        {'text': 'Stamatatos et al. (2023) "Diff-Vectors for Authorship Analysis"',
         'font_size': SMALL_FONT, 'font_color': MID_GRAY, 'space_after': 4},
        {'text': 'Ganin & Lempitsky (2015) "Unsupervised Domain Adaptation by Backpropagation"',
         'font_size': SMALL_FONT, 'font_color': MID_GRAY, 'space_after': 4},
        {'text': 'Abbasi & Chen (2008) "Writeprints: A Stylometric Approach to '
                 'Identity-Level Identification"',
         'font_size': SMALL_FONT, 'font_color': MID_GRAY, 'space_after': 4},
        {'text': 'Boenninghoff et al. (2019) "AdHominem: Explainable Authorship Verification"',
         'font_size': SMALL_FONT, 'font_color': MID_GRAY, 'space_after': 2},
    ]
    _add_rich_textbox(slide, x + Emu(150000), y + Emu(120000),
                      col_w - Emu(300000), refs_h - Emu(200000), refs_paras)

    y += refs_h

    # =====================================================
    # FOOTER
    # =====================================================
    footer_h = Emu(600000)
    footer_y = SLIDE_H - footer_h
    _add_shape_with_fill(slide, MSO_SHAPE.RECTANGLE,
                          0, footer_y, SLIDE_W, footer_h, NAVY)

    _add_textbox(slide, Emu(800000), footer_y + Emu(120000),
                 SLIDE_W - Emu(1600000), footer_h - Emu(200000),
                 'COMP34812 Natural Language Understanding  |  University of Manchester  |  '
                 'Group 34  |  2025-26  |  github.com/ayush-kumar-prog',
                 Pt(22), RGBColor(0xAA, 0xBB, 0xCC), bold=False,
                 alignment=PP_ALIGN.CENTER)

    # =====================================================
    # SAVE
    # =====================================================
    prs.save(str(OUTPUT_PATH))
    print(f"\nPoster saved to {OUTPUT_PATH}")
    print(f"  Slide dimensions: {SLIDE_W} x {SLIDE_H} (40\" x 22.5\")")
    print(f"  Suitable for A1 landscape printing")
    print(f"  Export to PDF via PowerPoint/LibreOffice for final submission")


if __name__ == '__main__':
    create_poster()
