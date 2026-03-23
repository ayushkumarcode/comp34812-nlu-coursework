#!/usr/bin/env python3
"""
COMP34812 — Academic Poster Generator
Group 34

Creates an A1 academic poster using python-pptx.
Single slide, 16:9 aspect ratio.
Generates matplotlib visualizations and embeds them.

Usage: python poster/create_poster.py
"""

import os
import sys
from pathlib import Path

# Use poster venv if available
VENV = Path('/tmp/poster_env')
if VENV.exists():
    sys.path.insert(0, str(VENV / 'lib' / 'python3.14' / 'site-packages'))

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import numpy as np

# ============================================================
# CONFIGURATION
# ============================================================

SLIDE_WIDTH = Inches(40)
SLIDE_HEIGHT = Inches(22.5)

TITLE_BG = RGBColor(0x1B, 0x3A, 0x5C)
HEADER_BG = RGBColor(0x2E, 0x86, 0xAB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x2D, 0x2D, 0x2D)

TITLE_SIZE = Pt(48)
HEADER_SIZE = Pt(28)
BODY_SIZE = Pt(16)


def add_section(slide, title, body, x, y, width, header_h=Emu(1500000)):
    """Add a section with colored header and body text."""
    header = slide.shapes.add_shape(1, x, y, width, header_h)
    header.fill.solid()
    header.fill.fore_color.rgb = HEADER_BG
    header.line.fill.background()

    txBox = slide.shapes.add_textbox(x + Emu(100000), y + Emu(150000),
                                      width - Emu(200000), header_h - Emu(300000))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = HEADER_SIZE
    p.font.color.rgb = WHITE
    p.font.bold = True

    body_y = y + header_h + Emu(100000)
    body_box = slide.shapes.add_textbox(x + Emu(100000), body_y,
                                         width - Emu(200000), Emu(15000000))
    bf = body_box.text_frame
    bf.word_wrap = True
    p = bf.paragraphs[0]
    p.text = body
    p.font.size = BODY_SIZE
    p.font.color.rgb = DARK


def create_poster(task='av'):
    """Create academic poster as PPTX."""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title banner
    shape = slide.shapes.add_shape(1, 0, 0, SLIDE_WIDTH, Emu(4000000))
    shape.fill.solid()
    shape.fill.fore_color.rgb = TITLE_BG
    shape.line.fill.background()

    txBox = slide.shapes.add_textbox(Emu(500000), Emu(500000),
                                      SLIDE_WIDTH - Emu(1000000), Emu(2500000))
    tf = txBox.text_frame
    tf.word_wrap = True
    task_name = 'Authorship Verification' if task == 'av' else 'NLI'
    p = tf.paragraphs[0]
    p.text = f'{task_name}: A Multi-Strategy Approach'
    p.font.size = TITLE_SIZE
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = 'COMP34812 — Group 34'
    p2.font.size = Pt(22)
    p2.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    p2.alignment = PP_ALIGN.CENTER

    # 3-column layout
    col_w = Emu(12000000)
    gap = Emu(200000)
    y_start = Emu(4500000)

    # Column 1
    x1 = Emu(200000)
    add_section(slide, 'Introduction & Dataset', _intro(task), x1, y_start, col_w)
    add_section(slide, 'Solution 1 (Cat A)', _sol1(task), x1, y_start + Emu(6000000), col_w)

    # Column 2
    x2 = x1 + col_w + gap
    add_section(slide, 'Solution 2', _sol2(task), x2, y_start, col_w)
    add_section(slide, 'Results', _results(task), x2, y_start + Emu(6000000), col_w)

    # Column 3
    x3 = x2 + col_w + gap
    add_section(slide, 'Error Analysis', _errors(task), x3, y_start, col_w)
    add_section(slide, 'Limitations & Ethics', _limits(task), x3, y_start + Emu(6000000), col_w)

    # Embed charts if they exist
    poster_dir = Path(__file__).parent
    for img_name in ['f1_chart.png', 'cm_av_cat_b.png']:
        img_path = poster_dir / img_name
        if img_path.exists():
            slide.shapes.add_picture(str(img_path), x2, y_start + Emu(9000000), Emu(11000000))

    prs.save('poster/poster.pptx')
    print("Poster saved to poster/poster.pptx")


def _intro(t):
    if t == 'nli':
        return ('NLI: determine if hypothesis is entailed by premise.\n'
                'Dataset: 24,432 train / 6,736 dev. Near-balanced.')
    return 'AV: determine if two texts share the same author.'


def _sol1(t):
    if t == 'nli':
        return 'Cat A — Stacking Ensemble, ~280 features.'
    return 'Cat A — ~950 stylometric features, stacking ensemble.'


def _sol2(t):
    if t == 'nli':
        return 'Cat C — DeBERTa-v3-base cross-encoder with GRL.'
    return 'Cat B — Siamese char-CNN + BiLSTM + GRL.'
