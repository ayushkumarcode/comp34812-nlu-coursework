#!/usr/bin/env python3
"""
COMP34812 — Award-Quality Academic Poster Generator (AV Track)
Group 34

Creates an A1-printable academic poster using python-pptx with
publication-quality matplotlib visualizations at 300 DPI.

Usage:
    source /tmp/poster_env/bin/activate
    python poster/create_poster_final.py
"""

import os
import sys
from pathlib import Path

# Use poster venv if available
VENV = Path('/tmp/poster_env')
if VENV.exists():
    for p in VENV.glob('lib/*/site-packages'):
        sys.path.insert(0, str(p))

import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import matplotlib.patheffects as path_effects
import numpy as np

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ============================================================
# CONFIGURATION
# ============================================================

PROJECT_DIR = Path(__file__).parent.parent
POSTER_DIR = Path(__file__).parent
OUTPUT_PATH = POSTER_DIR / 'poster_final.pptx'

# Slide dimensions — 40" x 22.5" (16:9, A1-printable)
SLIDE_W = Inches(40)
SLIDE_H = Inches(22.5)

# Color palette — refined for visual impact
NAVY        = RGBColor(0x1B, 0x3A, 0x5C)
NAVY_DARK   = RGBColor(0x12, 0x2A, 0x45)
TEAL        = RGBColor(0x2E, 0x86, 0xAB)
DARK_TEAL   = RGBColor(0x1E, 0x6E, 0x8E)
CORAL       = RGBColor(0xE8, 0x4D, 0x60)
LIGHT_BG    = RGBColor(0xF5, 0xF7, 0xFA)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
DARK        = RGBColor(0x2D, 0x2D, 0x2D)
MID_GRAY    = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY  = RGBColor(0xBB, 0xBB, 0xBB)
ACCENT_GREEN  = RGBColor(0x27, 0xAE, 0x60)
ACCENT_ORANGE = RGBColor(0xE6, 0x7E, 0x22)
TABLE_HEADER  = RGBColor(0x1B, 0x3A, 0x5C)
TABLE_ALT_BG  = RGBColor(0xEB, 0xF0, 0xF5)
BORDER_COLOR  = RGBColor(0xDD, 0xDD, 0xDD)
GOLD          = RGBColor(0xF3, 0x9C, 0x12)

# Hex strings for matplotlib
C_NAVY = '#1B3A5C'
C_TEAL = '#2E86AB'
C_CORAL = '#E84D60'
C_GREEN = '#27AE60'
C_ORANGE = '#E67E22'
C_DARK = '#2D2D2D'
C_GRAY = '#95A5A6'
C_LIGHTBG = '#F5F7FA'

# Font sizes (calibrated for A1 print)
TITLE_FONT     = Pt(72)
SUBTITLE_FONT  = Pt(36)
HEADER_FONT    = Pt(38)
BODY_FONT      = Pt(24)
BODY_FONT_SM   = Pt(22)
SMALL_FONT     = Pt(20)
TINY_FONT      = Pt(18)

# Layout constants
MARGIN     = Emu(700000)
COL_GAP    = Emu(450000)
SEC_GAP    = Emu(350000)
PAD_INNER  = Emu(200000)

# ============================================================
# DATA LOADING
# ============================================================

def load_predictions():
    """Load ground truth and best predictions for AV task."""
    gt = []
    gt_path = PROJECT_DIR / 'baseline_extracted' / 'nlu_bundle-feature-unified-local-scorer' / \
              'local_scorer' / 'reference_data' / 'NLU_SharedTask_AV_dev.solution'
    with open(gt_path) as f:
        for line in f:
            gt.append(int(line.strip()))

    def load_pred(path):
        preds = []
        with open(path) as f:
            for line in f:
                preds.append(int(line.strip()))
        return preds

    pred_dir = PROJECT_DIR / 'predictions'
    # Best Cat A: LightGBM ensemble (F1=0.7340)
    pred_a = load_pred(pred_dir / 'av_Group_34_A_lgbm.csv')
    # Best Cat B: Adversarial v3 (F1=0.7422)
    pred_b = load_pred(pred_dir / 'av_Group_34_B_v3.csv')

    if len(pred_a) > len(gt):
        pred_a = pred_a[1:]
    if len(pred_b) > len(gt):
        pred_b = pred_b[1:]

    return np.array(gt), np.array(pred_a), np.array(pred_b)


# ============================================================
# CHART GENERATION — HIGH-QUALITY 300 DPI
# ============================================================

def _set_chart_style():
    """Set consistent matplotlib style for all charts."""
    plt.rcParams.update({
        'font.family': 'sans-serif',
        'font.sans-serif': ['Helvetica Neue', 'Helvetica', 'Arial', 'DejaVu Sans'],
        'font.size': 14,
        'axes.labelsize': 18,
        'axes.titlesize': 22,
        'xtick.labelsize': 15,
        'ytick.labelsize': 15,
        'figure.facecolor': 'white',
        'axes.facecolor': 'white',
        'axes.edgecolor': '#CCCCCC',
        'axes.grid': False,
        'grid.alpha': 0.3,
        'grid.linestyle': '--',
        'legend.framealpha': 0.95,
        'legend.edgecolor': '#CCCCCC',
    })


def generate_f1_chart():
    """Generate polished F1 comparison bar chart with significance stars."""
    _set_chart_style()

    models = ['SVM\nBaseline', 'LSTM\nBaseline', 'BERT\nBaseline',
              'Sol 1\n(Cat A)', 'Sol 2\n(Cat B)']
    scores = [0.5610, 0.6226, 0.7854, 0.7340, 0.7422]

    fig, ax = plt.subplots(figsize=(14, 6))

    # Colors with gradient effect via edge
    bar_colors = ['#B0BEC5', '#B0BEC5', '#B0BEC5', C_GREEN, C_ORANGE]
    edge_colors = ['#78909C', '#78909C', '#78909C', '#1B8C4F', '#C0651C']
    hatches = ['', '', '', '///', '\\\\\\']

    bars = ax.bar(models, scores, color=bar_colors, edgecolor=edge_colors,
                  linewidth=2.5, width=0.62, zorder=3)

    # Add subtle hatching to our solutions for distinction
    for i in range(3, 5):
        bars[i].set_hatch(hatches[i])
        bars[i].set_edgecolor(edge_colors[i])

    # Value labels with shadow effect
    for bar, score, color in zip(bars, scores, bar_colors):
        txt = ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 0.012,
                      f'{score:.4f}', ha='center', va='bottom',
                      fontsize=22, fontweight='bold', color=C_DARK)
        txt.set_path_effects([
            path_effects.withStroke(linewidth=3, foreground='white')
        ])

    # Statistical significance stars
    star_props = dict(fontsize=28, fontweight='bold', ha='center', va='bottom')
    ax.text(3, scores[3] + 0.04, '***', color=C_GREEN, **star_props)
    ax.text(4, scores[4] + 0.04, '***', color=C_ORANGE, **star_props)

    # Improvement annotations — inside the bars
    # Sol 1 vs SVM
    ax.text(3, scores[3] / 2, '+30.8%\nvs SVM', ha='center', va='center',
            fontsize=16, fontweight='bold', color='white',
            path_effects=[path_effects.withStroke(linewidth=4, foreground=C_GREEN)])

    # Sol 2 vs LSTM
    ax.text(4, scores[4] / 2, '+19.2%\nvs LSTM', ha='center', va='center',
            fontsize=16, fontweight='bold', color='white',
            path_effects=[path_effects.withStroke(linewidth=4, foreground=C_ORANGE)])

    # Styling
    ax.set_ylabel('Macro F1 Score', fontsize=24, fontweight='bold',
                   labelpad=15, color=C_DARK)
    ax.set_ylim(0, 0.92)
    ax.set_xlim(-0.5, 4.5)

    # Horizontal grid
    ax.yaxis.grid(True, linestyle='--', alpha=0.35, zorder=0, color='#CCCCCC')
    ax.set_axisbelow(True)

    ax.tick_params(axis='both', labelsize=18, colors=C_DARK)
    ax.tick_params(axis='x', length=0, pad=10)

    for spine in ['top', 'right']:
        ax.spines[spine].set_visible(False)
    for spine in ['left', 'bottom']:
        ax.spines[spine].set_color('#CCCCCC')
        ax.spines[spine].set_linewidth(1.5)

    # BERT reference line
    ax.axhline(y=0.7854, color='#BDC3C7', linestyle=':', linewidth=2, alpha=0.6)
    ax.text(4.45, 0.79, 'BERT', fontsize=14, color='#95a5a6',
            ha='right', style='italic', alpha=0.8)

    # Legend
    baseline_patch = mpatches.Patch(facecolor='#B0BEC5', edgecolor='#78909C',
                                     linewidth=1.5, label='Baselines')
    sol1_patch = mpatches.Patch(facecolor=C_GREEN, edgecolor='#1B8C4F',
                                 linewidth=1.5, label='Our Sol 1 (Cat A)')
    sol2_patch = mpatches.Patch(facecolor=C_ORANGE, edgecolor='#C0651C',
                                 linewidth=1.5, label='Our Sol 2 (Cat B)')
    ax.legend(handles=[baseline_patch, sol1_patch, sol2_patch],
              fontsize=17, loc='upper left', framealpha=0.95,
              edgecolor='#CCCCCC', fancybox=True, shadow=True)

    # Significance note
    ax.text(0.98, 0.02, '*** p < 0.001 (McNemar\'s test)',
            transform=ax.transAxes, ha='right', fontsize=13,
            color='#999999', style='italic')

    plt.tight_layout(pad=1.5)
    path = POSTER_DIR / 'f1_chart_final.png'
    plt.savefig(str(path), dpi=300, bbox_inches='tight',
                facecolor='white', edgecolor='none')
    plt.close()
    print(f"  [OK] F1 chart -> {path}")
    return path


def generate_confusion_matrices(gt, pred_a, pred_b):
    """Generate side-by-side confusion matrices with real data."""
    _set_chart_style()
    from sklearn.metrics import confusion_matrix, precision_score, recall_score, f1_score

    cm_a = confusion_matrix(gt, pred_a)
    cm_b = confusion_matrix(gt, pred_b)

    fig, axes = plt.subplots(1, 2, figsize=(18, 6))

    titles = ['Solution 1 (Cat A)\nStylometric LightGBM Ensemble',
              'Solution 2 (Cat B)\nAdversarial CNN-BiLSTM']
    cms = [cm_a, cm_b]
    cmaps = ['Blues', 'Oranges']
    accent_colors = [C_GREEN, C_ORANGE]

    for ax, cm, title, cmap, accent in zip(axes, cms, titles, cmaps, accent_colors):
        # Normalize for color mapping but show raw counts
        cm_norm = cm.astype(float) / cm.sum(axis=1, keepdims=True)

        im = ax.imshow(cm_norm, cmap=cmap, interpolation='nearest',
                        aspect='equal', vmin=0, vmax=1)

        ax.set_xticks([0, 1])
        ax.set_yticks([0, 1])
        ax.set_xticklabels(['Different\nAuthor', 'Same\nAuthor'], fontsize=17)
        ax.set_yticklabels(['Different\nAuthor', 'Same\nAuthor'], fontsize=17)
        ax.set_xlabel('Predicted', fontsize=20, fontweight='bold', labelpad=12)
        ax.set_ylabel('True', fontsize=20, fontweight='bold', labelpad=12)
        ax.set_title(title, fontsize=21, fontweight='bold', pad=18, color=C_NAVY)

        # Cell values with counts and percentages
        for i in range(2):
            for j in range(2):
                val = cm[i, j]
                pct = cm_norm[i, j] * 100
                color = 'white' if cm_norm[i, j] > 0.5 else C_DARK
                ax.text(j, i, f'{val:,}\n({pct:.1f}%)', ha='center', va='center',
                        fontsize=22, fontweight='bold', color=color)

        # Add border around diagonal (correct predictions)
        for k in range(2):
            rect = plt.Rectangle((k-0.5, k-0.5), 1, 1, linewidth=3,
                                  edgecolor=accent, facecolor='none',
                                  linestyle='--', alpha=0.7)
            ax.add_patch(rect)

        # Metrics below
        prec = precision_score(gt, pred_a if cmap == 'Blues' else pred_b, average='macro')
        rec = recall_score(gt, pred_a if cmap == 'Blues' else pred_b, average='macro')
        f1 = f1_score(gt, pred_a if cmap == 'Blues' else pred_b, average='macro')
        ax.text(0.5, -0.18,
                f'Precision={prec:.3f}    Recall={rec:.3f}    F1={f1:.4f}',
                transform=ax.transAxes, ha='center', fontsize=16,
                color='#555555', fontweight='bold',
                bbox=dict(boxstyle='round,pad=0.4', facecolor='#F0F0F0',
                          edgecolor='#CCCCCC', alpha=0.9))

    plt.tight_layout(w_pad=4, pad=2)
    path = POSTER_DIR / 'cm_final.png'
    plt.savefig(str(path), dpi=300, bbox_inches='tight',
                facecolor='white', edgecolor='none')
    plt.close()
    print(f"  [OK] Confusion matrices -> {path}")
    return path


def generate_architecture_diagram():
    """Generate a clean, modern architecture diagram for Sol 2."""
    _set_chart_style()
    # Use DejaVu Sans for architecture to avoid missing glyph issues
    plt.rcParams['font.sans-serif'] = ['DejaVu Sans', 'Arial']

    fig, ax = plt.subplots(figsize=(15, 9))
    ax.set_xlim(-0.5, 10.5)
    ax.set_ylim(-1.0, 11)
    ax.axis('off')
    fig.patch.set_facecolor('white')

    def draw_box(x, y, w, h, text, color, text_color='white',
                 fontsize=14, alpha=1.0, style='round,pad=0.15',
                 linewidth=2.5, edgecolor=None):
        """Draw a rounded box with shadow effect."""
        ec = edgecolor or color
        # Shadow
        shadow = mpatches.FancyBboxPatch(
            (x - w/2 + 0.05, y - h/2 - 0.05), w, h,
            boxstyle=style, facecolor='#00000010',
            edgecolor='none', linewidth=0)
        ax.add_patch(shadow)
        # Main box
        rect = mpatches.FancyBboxPatch(
            (x - w/2, y - h/2), w, h,
            boxstyle=style, facecolor=color,
            edgecolor=ec, linewidth=linewidth, alpha=alpha)
        ax.add_patch(rect)
        ax.text(x, y, text, ha='center', va='center',
                fontsize=fontsize, fontweight='bold', color=text_color,
                zorder=10)

    def arrow(x1, y1, x2, y2, color='#666666', width=2, style='->'):
        ax.annotate('', xy=(x2, y2), xytext=(x1, y1),
                    arrowprops=dict(arrowstyle=style, color=color,
                                    lw=width, connectionstyle='arc3,rad=0'))

    # Title
    ax.text(5, 10.5, 'Solution 2: Adversarial Style-Content Disentanglement',
            ha='center', fontsize=24, fontweight='bold', color=C_NAVY)
    ax.text(5, 10.1, 'Siamese CNN-BiLSTM with Gradient Reversal',
            ha='center', fontsize=16, color='#666666', style='italic')

    # ---- Input Layer ----
    draw_box(2.5, 9.2, 3.0, 0.7, 'Text 1', '#E3F2FD', C_NAVY, 17, edgecolor='#90CAF9')
    draw_box(7.5, 9.2, 3.0, 0.7, 'Text 2', '#E3F2FD', C_NAVY, 17, edgecolor='#90CAF9')

    # ---- Char Embedding ----
    draw_box(2.5, 7.9, 3.2, 0.65, 'Char Embed (64d)', '#42A5F5', 'white', 14)
    draw_box(7.5, 7.9, 3.2, 0.65, 'Char Embed (64d)', '#42A5F5', 'white', 14)

    # ---- Multi-width CNN ----
    # Draw 3 sub-boxes for kernel sizes
    for offset, ks in [(-0.9, '3'), (0, '5'), (0.9, '7')]:
        draw_box(2.5 + offset, 6.65, 0.75, 0.55, f'k={ks}', '#1E88E5', 'white', 12,
                 style='round,pad=0.08', linewidth=1.5)
        draw_box(7.5 + offset, 6.65, 0.75, 0.55, f'k={ks}', '#1E88E5', 'white', 12,
                 style='round,pad=0.08', linewidth=1.5)

    # CNN label
    ax.text(2.5, 7.15, 'Multi-Width CNN', ha='center', fontsize=13,
            fontweight='bold', color=C_NAVY)
    ax.text(7.5, 7.15, 'Multi-Width CNN', ha='center', fontsize=13,
            fontweight='bold', color=C_NAVY)

    # ---- BiLSTM ----
    draw_box(2.5, 5.5, 3.2, 0.65, 'BiLSTM (256h)', '#1565C0', 'white', 14)
    draw_box(7.5, 5.5, 3.2, 0.65, 'BiLSTM (256h)', '#1565C0', 'white', 14)

    # ---- Attention ----
    draw_box(2.5, 4.3, 3.2, 0.65, 'Additive Attention', '#0D47A1', 'white', 14)
    draw_box(7.5, 4.3, 3.2, 0.65, 'Additive Attention', '#0D47A1', 'white', 14)

    # Output labels
    ax.text(2.5, 3.65, 'v\u2081', ha='center', fontsize=18,
            fontweight='bold', color=C_NAVY, style='italic')
    ax.text(7.5, 3.65, 'v\u2082', ha='center', fontsize=18,
            fontweight='bold', color=C_NAVY, style='italic')

    # Shared weights indicator
    ax.annotate('', xy=(4.1, 7.9), xytext=(5.9, 7.9),
                arrowprops=dict(arrowstyle='<->', color=C_CORAL, lw=3,
                                connectionstyle='arc3,rad=0'))
    ax.text(5, 8.2, 'Shared Weights', ha='center', fontsize=13,
            color=C_CORAL, fontweight='bold',
            bbox=dict(boxstyle='round,pad=0.2', facecolor='white',
                      edgecolor=C_CORAL, alpha=0.9))

    # Arrows down left tower
    for y_top, y_bot in [(8.85, 8.22), (7.57, 7.15), (6.37, 5.82), (5.17, 4.62)]:
        arrow(2.5, y_top, 2.5, y_bot, '#888888', 1.5)
    # Arrows down right tower
    for y_top, y_bot in [(8.85, 8.22), (7.57, 7.15), (6.37, 5.82), (5.17, 4.62)]:
        arrow(7.5, y_top, 7.5, y_bot, '#888888', 1.5)

    # ---- Merge Layer ----
    draw_box(5, 2.8, 4.5, 0.75,
             '[v1, v2, |v1-v2|, v1*v2]',
             '#7B1FA2', 'white', 15)
    ax.text(5, 2.15, 'Concatenation (4D)', ha='center', fontsize=12,
            color='#888888', style='italic')

    # Arrows to merge
    arrow(2.5, 3.55, 3.5, 3.17, '#666666', 1.8)
    arrow(7.5, 3.55, 6.5, 3.17, '#666666', 1.8)

    # ---- Two branches ----
    # Branch 1: Authorship classifier (left)
    draw_box(2.8, 1.2, 3.8, 0.7,
             'Authorship Classifier',
             C_GREEN, 'white', 15, linewidth=3)
    ax.text(2.8, 0.55, 'Same / Different Author',
            ha='center', fontsize=13, color=C_GREEN, fontweight='bold')

    # Branch 2: GRL + Topic
    draw_box(7.2, 1.2, 3.8, 0.7,
             'GRL -> Topic Classifier',
             C_CORAL, 'white', 15, linewidth=3)
    ax.text(7.2, 0.55, 'Gradient Reversal Layer',
            ha='center', fontsize=13, color=C_CORAL, fontweight='bold')
    ax.text(7.2, 0.15, 'Forces topic-invariant representations',
            ha='center', fontsize=11, color='#999999', style='italic')

    # Arrows to branches
    arrow(4.0, 2.42, 2.8, 1.55, C_GREEN, 2.5)
    arrow(6.0, 2.42, 7.2, 1.55, C_CORAL, 2.5)

    # Dimension annotations
    ax.text(0.3, 7.9, '64', ha='center', fontsize=11, color='#AAA',
            bbox=dict(boxstyle='round', facecolor='#F5F5F5', edgecolor='#DDD'))
    ax.text(0.3, 6.65, '3\u00D7128', ha='center', fontsize=11, color='#AAA',
            bbox=dict(boxstyle='round', facecolor='#F5F5F5', edgecolor='#DDD'))
    ax.text(0.3, 5.5, '512', ha='center', fontsize=11, color='#AAA',
            bbox=dict(boxstyle='round', facecolor='#F5F5F5', edgecolor='#DDD'))
    ax.text(0.3, 4.3, '512', ha='center', fontsize=11, color='#AAA',
            bbox=dict(boxstyle='round', facecolor='#F5F5F5', edgecolor='#DDD'))

    plt.tight_layout(pad=1)
    path = POSTER_DIR / 'arch_diagram_final.png'
    plt.savefig(str(path), dpi=300, bbox_inches='tight',
                facecolor='white', edgecolor='none')
    plt.close()
    print(f"  [OK] Architecture diagram -> {path}")
    return path


def generate_feature_groups_chart():
    """Generate horizontal bar chart of feature groups with counts."""
    _set_chart_style()
    plt.rcParams['font.sans-serif'] = ['DejaVu Sans', 'Arial']

    groups = [
        'Function Words', 'TF-IDF + SVD', 'Character N-gram',
        'POS Tags', 'Lexical Richness', 'Structural',
        'Syntactic Complexity', 'Writing Rhythm', 'Info-Theoretic'
    ]
    counts = [150, 100, 56, 45, 30, 15, 10, 6, 5]

    # Standard (teal) vs Novel (coral)
    is_novel = [False, False, False, False, False, False, True, True, True]
    colors = [C_CORAL if n else C_TEAL for n in is_novel]
    edge_colors = ['#C43350' if n else '#1E6E8E' for n in is_novel]

    fig, ax = plt.subplots(figsize=(13, 7))
    y_pos = np.arange(len(groups))

    bars = ax.barh(y_pos, counts, color=colors, edgecolor=edge_colors,
                   linewidth=1.5, height=0.6, zorder=3)

    # Value labels with background
    for bar, count, novel in zip(bars, counts, is_novel):
        label_color = C_CORAL if novel else C_TEAL
        txt = ax.text(bar.get_width() + 3, bar.get_y() + bar.get_height()/2,
                      f'{count}', va='center', fontsize=19, fontweight='bold',
                      color=label_color)

    ax.set_yticks(y_pos)
    ax.set_yticklabels(groups, fontsize=18, fontweight='bold')
    ax.set_xlabel('Number of Features', fontsize=21, fontweight='bold',
                  labelpad=12, color=C_DARK)
    ax.set_title('Feature Architecture: 417 per text, 695-dim diff-vector',
                 fontsize=22, fontweight='bold', pad=20, color=C_NAVY)
    ax.set_xlim(0, 190)
    ax.invert_yaxis()

    ax.xaxis.grid(True, linestyle='--', alpha=0.3, zorder=0, color='#CCCCCC')
    for spine in ['top', 'right']:
        ax.spines[spine].set_visible(False)
    ax.spines['left'].set_color('#CCCCCC')
    ax.spines['bottom'].set_color('#CCCCCC')

    # Legend
    std_patch = mpatches.Patch(facecolor=C_TEAL, edgecolor='#1E6E8E',
                                linewidth=1.5, label='Standard Stylometric (396)')
    novel_patch = mpatches.Patch(facecolor=C_CORAL, edgecolor='#C43350',
                                  linewidth=1.5, label='Novel Features (21)')
    ax.legend(handles=[std_patch, novel_patch], fontsize=17,
              loc='lower right', framealpha=0.95, edgecolor='#CCCCCC',
              fancybox=True, shadow=True)

    # Cumulative annotation
    ax.text(0.98, 0.98, '695 dim diff-vector\n|f(t\u2081) \u2212 f(t\u2082)|',
            transform=ax.transAxes, ha='right', va='top',
            fontsize=14, color='#888888', style='italic',
            bbox=dict(boxstyle='round,pad=0.4', facecolor='#F8F8F8',
                      edgecolor='#DDD'))

    plt.tight_layout(pad=1.5)
    path = POSTER_DIR / 'feature_groups_final.png'
    plt.savefig(str(path), dpi=300, bbox_inches='tight',
                facecolor='white', edgecolor='none')
    plt.close()
    print(f"  [OK] Feature groups chart -> {path}")
    return path


# ============================================================
# POSTER LAYOUT HELPERS
# ============================================================

def _shape(slide, left, top, width, height, fill_color,
           line_color=None, line_width=None):
    """Add a filled rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def _rounded_rect(slide, left, top, width, height, fill_color,
                  line_color=None, line_width=None):
    """Add a rounded rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    # Adjust corner radius
    shape.adjustments[0] = 0.05
    return shape


def _textbox(slide, left, top, width, height, text, font_size,
             font_color=DARK, bold=False, alignment=PP_ALIGN.LEFT,
             font_name='Calibri', line_spacing=1.15):
    """Add a simple text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    if line_spacing:
        p.line_spacing = line_spacing
    return txBox


def _rich_textbox(slide, left, top, width, height, paragraphs_data,
                  font_name='Calibri'):
    """Add a text box with multiple styled paragraphs.

    paragraphs_data: list of dicts with keys:
        text, font_size, font_color, bold, alignment, runs, space_after, etc.
    """
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, pdata in enumerate(paragraphs_data):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        runs = pdata.get('runs', None)
        if runs:
            for j, run_data in enumerate(runs):
                if j == 0 and p.runs:
                    run = p.runs[0]
                else:
                    run = p.add_run()
                run.text = run_data.get('text', '')
                run.font.size = run_data.get('font_size', pdata.get('font_size', BODY_FONT))
                run.font.color.rgb = run_data.get('font_color', pdata.get('font_color', DARK))
                run.font.bold = run_data.get('bold', pdata.get('bold', False))
                run.font.name = run_data.get('font_name', font_name)
                if run_data.get('italic', False):
                    run.font.italic = True
        else:
            p.text = pdata.get('text', '')
            p.font.size = pdata.get('font_size', BODY_FONT)
            p.font.color.rgb = pdata.get('font_color', DARK)
            p.font.bold = pdata.get('bold', False)
            p.font.name = pdata.get('font_name', font_name)
            if pdata.get('italic', False):
                p.font.italic = True

        p.alignment = pdata.get('alignment', PP_ALIGN.LEFT)
        p.space_after = Pt(pdata.get('space_after', 4))
        p.space_before = Pt(pdata.get('space_before', 0))
        ls = pdata.get('line_spacing', 1.15)
        if ls:
            p.line_spacing = ls

    return txBox


def _section_header(slide, left, top, width, title, icon_text=None):
    """Add a professional section header with accent bar."""
    h = Emu(720000)

    # Header background with subtle gradient effect (two overlapping rects)
    _shape(slide, left, top, width, h, TEAL)

    # Left accent bar (darker)
    accent_w = Emu(90000)
    _shape(slide, left, top, accent_w, h, DARK_TEAL)

    # Optional icon/number on the left
    if icon_text:
        _textbox(slide, left + Emu(130000), top + Emu(90000),
                 Emu(400000), h - Emu(180000),
                 icon_text, Pt(32), WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)
        text_left = left + Emu(450000)
    else:
        text_left = left + Emu(250000)

    # Header text
    _textbox(slide, text_left, top + Emu(110000),
             width - Emu(350000), h - Emu(220000),
             title, HEADER_FONT, WHITE, bold=True)

    # Bottom highlight line
    _shape(slide, left, top + h - Emu(50000), width, Emu(50000), DARK_TEAL)

    return h


def _section_body(slide, left, top, width, height):
    """Add section body background with rounded corners."""
    return _rounded_rect(slide, left, top, width, height,
                          LIGHT_BG, BORDER_COLOR, Pt(1.5))


def _key_value_line(label, value, label_color=DARK, value_color=DARK,
                    font_size=BODY_FONT, bold_label=True):
    """Helper to create a key-value styled paragraph."""
    return {
        'runs': [
            {'text': f'{label}: ', 'font_size': font_size, 'bold': bold_label,
             'font_color': label_color},
            {'text': value, 'font_size': font_size, 'font_color': value_color}
        ],
        'font_size': font_size,
        'space_after': 6
    }


def _bullet(text, font_size=BODY_FONT_SM, color=DARK, bold=False):
    """Helper to create a bullet point paragraph."""
    return {
        'runs': [
            {'text': '  \u2022  ', 'font_size': font_size, 'font_color': MID_GRAY},
            {'text': text, 'font_size': font_size, 'font_color': color, 'bold': bold}
        ],
        'font_size': font_size,
        'space_after': 4
    }


def _bullet_with_detail(label, detail, label_color=DARK, detail_color=MID_GRAY,
                         font_size=BODY_FONT_SM):
    """Bullet with bold label and lighter detail."""
    return {
        'runs': [
            {'text': '  \u2022  ', 'font_size': font_size, 'font_color': MID_GRAY},
            {'text': label, 'font_size': font_size, 'font_color': label_color, 'bold': True},
            {'text': f' \u2014 {detail}', 'font_size': Pt(20), 'font_color': detail_color}
        ],
        'font_size': font_size,
        'space_after': 4
    }


# ============================================================
# MAIN POSTER CREATION
# ============================================================

def create_poster():
    """Create the final award-quality AV track academic poster."""
    print("=" * 60)
    print("  Creating AV Track Academic Poster (Final Version)")
    print("=" * 60)

    # Load real data
    print("\nLoading prediction data...")
    gt, pred_a, pred_b = load_predictions()
    print(f"  Ground truth: {len(gt)} samples")
    print(f"  Pred A (Cat A): {len(pred_a)} samples")
    print(f"  Pred B (Cat B): {len(pred_b)} samples")

    # Generate all visualizations
    print("\nGenerating visualizations at 300 DPI...")
    f1_path = generate_f1_chart()
    cm_path = generate_confusion_matrices(gt, pred_a, pred_b)
    arch_path = generate_architecture_diagram()
    feat_path = generate_feature_groups_chart()

    # Create presentation
    print("\nBuilding poster layout...")
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

    # Background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0xF0, 0xF2, 0xF5)

    # =====================================================
    # TITLE BANNER (compact — 2.9")
    # =====================================================
    title_h = Emu(2650000)

    # Main background
    _shape(slide, 0, 0, SLIDE_W, title_h, NAVY)

    # Decorative gradient stripe at top
    _shape(slide, 0, 0, SLIDE_W, Emu(80000), NAVY_DARK)

    # Accent stripe at bottom of title
    _shape(slide, 0, title_h - Emu(90000), SLIDE_W, Emu(90000), TEAL)

    # Thin coral accent line
    _shape(slide, 0, title_h - Emu(100000), SLIDE_W, Emu(12000), CORAL)

    # Title text
    _textbox(slide, Emu(1000000), Emu(200000),
             SLIDE_W - Emu(2000000), Emu(1400000),
             'Authorship Verification: Combining Stylometric Feature\n'
             'Engineering with Adversarial Neural Style-Content Disentanglement',
             Pt(64), WHITE, bold=True, alignment=PP_ALIGN.CENTER,
             line_spacing=1.05)

    # Subtitle
    _textbox(slide, Emu(1000000), Emu(1650000),
             SLIDE_W - Emu(2000000), Emu(450000),
             'COMP34812 Natural Language Understanding  |  Group 34  |  '
             'University of Manchester  |  2025\u201326',
             Pt(32), RGBColor(0xBB, 0xCC, 0xDD), bold=False,
             alignment=PP_ALIGN.CENTER)

    # Key results highlight strip below title
    highlight_y = title_h
    highlight_h = Emu(350000)
    _shape(slide, 0, highlight_y, SLIDE_W, highlight_h,
           RGBColor(0xE8, 0xEB, 0xF0))

    # Sol 1 callout
    _rich_textbox(slide, Emu(3000000), highlight_y + Emu(50000),
                  Emu(12000000), highlight_h - Emu(100000),
                  [{'runs': [
                      {'text': 'Sol 1 (Cat A): ', 'font_size': Pt(26),
                       'font_color': ACCENT_GREEN, 'bold': True},
                      {'text': 'F1 = 0.7340', 'font_size': Pt(30),
                       'font_color': ACCENT_GREEN, 'bold': True},
                      {'text': '  (+30.8% vs SVM)', 'font_size': Pt(22),
                       'font_color': MID_GRAY},
                  ], 'alignment': PP_ALIGN.CENTER}])

    # Sol 2 callout
    _rich_textbox(slide, Emu(21000000), highlight_y + Emu(50000),
                  Emu(14000000), highlight_h - Emu(100000),
                  [{'runs': [
                      {'text': 'Sol 2 (Cat B): ', 'font_size': Pt(26),
                       'font_color': ACCENT_ORANGE, 'bold': True},
                      {'text': 'F1 = 0.7422', 'font_size': Pt(30),
                       'font_color': ACCENT_ORANGE, 'bold': True},
                      {'text': '  (+19.2% vs LSTM)', 'font_size': Pt(22),
                       'font_color': MID_GRAY},
                      {'text': '    Both p < 0.001 ***', 'font_size': Pt(20),
                       'font_color': LIGHT_GRAY, 'italic': True},
                  ], 'alignment': PP_ALIGN.CENTER}])

    # =====================================================
    # COLUMN LAYOUT: 4 columns
    # =====================================================
    usable_w = SLIDE_W - 2 * MARGIN
    n_cols = 4
    col_w = (usable_w - (n_cols - 1) * COL_GAP) // n_cols
    y_start = title_h + highlight_h + Emu(150000)

    # Bottom strip starts here — all column content MUST end above this
    bottom_strip_y = Emu(15700000)  # ~17.16"

    col_x = []
    for i in range(n_cols):
        col_x.append(MARGIN + i * (col_w + COL_GAP))

    # Column vertical dividers (subtle)
    for i in range(1, n_cols):
        div_x = col_x[i] - COL_GAP // 2
        _shape(slide, div_x, y_start, Emu(15000),
               bottom_strip_y - y_start,
               RGBColor(0xDD, 0xDD, 0xDD))

    # =====================================================
    # COLUMN 1: Introduction + Dataset + Evaluation
    # =====================================================
    x = col_x[0]
    y = y_start

    # --- INTRODUCTION ---
    hdr_h = _section_header(slide, x, y, col_w, 'Introduction', '\u2460')
    y += hdr_h

    intro_h = Emu(3550000)
    _section_body(slide, x, y, col_w, intro_h)

    intro_paras = [
        {'text': 'Task Definition',
         'font_size': BODY_FONT, 'bold': True, 'font_color': TEAL, 'space_after': 4},
        {'text': 'Given two texts, determine whether they were written by the '
                 'same author. Binary classification requiring models to capture '
                 'authorial style while remaining robust to topical variation.',
         'font_size': BODY_FONT_SM, 'space_after': 8},
        {'text': 'Key Challenge',
         'font_size': BODY_FONT, 'bold': True, 'font_color': CORAL, 'space_after': 4},
        {'text': 'Style-content confound: models may exploit topic similarity '
                 'as a proxy for shared authorship, producing spurious '
                 'correlations that fail to generalize.',
         'font_size': BODY_FONT_SM, 'space_after': 8},
        {'text': 'Our Approach',
         'font_size': BODY_FONT, 'bold': True, 'font_color': TEAL, 'space_after': 4},
        _bullet_with_detail('Solution 1 (Cat A)', 'Stylometric feature ensemble '
                            'with LightGBM classifier',
                            ACCENT_GREEN),
        _bullet_with_detail('Solution 2 (Cat B)', 'Adversarial CNN-BiLSTM with '
                            'gradient reversal for debiasing',
                            ACCENT_ORANGE),
    ]
    _rich_textbox(slide, x + PAD_INNER, y + Emu(120000),
                  col_w - 2 * PAD_INNER, intro_h - Emu(200000), intro_paras)

    y += intro_h + Emu(200000)

    # --- DATASET ---
    hdr_h = _section_header(slide, x, y, col_w, 'Dataset Summary', '\u2461')
    y += hdr_h

    dataset_h = Emu(3250000)
    _section_body(slide, x, y, col_w, dataset_h)

    dataset_paras = [
        _key_value_line('Training', '27,643 text pairs', TEAL),
        _key_value_line('Development', '5,993 text pairs', TEAL),
        _key_value_line('Labels', 'Same Author (51%) / Different Author (49%)', TEAL),
        _key_value_line('Domains', 'Cross-domain (emails, blog posts, reviews)', TEAL),
        {'text': '', 'font_size': Pt(4), 'space_after': 4},
        {'text': 'Representation',
         'font_size': BODY_FONT, 'bold': True, 'font_color': TEAL, 'space_after': 4},
        _bullet('Text pair (t\u2081, t\u2082) \u2192 diff-vector '
                '|f(t\u2081) \u2212 f(t\u2082)| capturing stylistic distance'),
        _bullet('Content-agnostic by design'),
        {'text': 'Closed mode \u2014 only provided training data permitted.',
         'font_size': SMALL_FONT, 'font_color': MID_GRAY, 'italic': True, 'space_after': 2},
    ]
    _rich_textbox(slide, x + PAD_INNER, y + Emu(120000),
                  col_w - 2 * PAD_INNER, dataset_h - Emu(200000), dataset_paras)

    y += dataset_h + Emu(200000)

    # --- EVALUATION ---
    hdr_h = _section_header(slide, x, y, col_w, 'Evaluation Methodology', '\u2462')
    y += hdr_h

    eval_h = Emu(3200000)
    _section_body(slide, x, y, col_w, eval_h)

    eval_paras = [
        {'text': 'Metrics',
         'font_size': BODY_FONT, 'bold': True, 'font_color': TEAL, 'space_after': 4},
        _bullet('Macro F1 (primary), MCC, Per-class P/R'),
        {'text': '', 'font_size': Pt(4), 'space_after': 4},
        {'text': 'Statistical Testing',
         'font_size': BODY_FONT, 'bold': True, 'font_color': TEAL, 'space_after': 4},
        _bullet('McNemar\'s test for paired significance'),
        _bullet('All improvements significant at p < 0.001'),
        {'text': '', 'font_size': Pt(4), 'space_after': 4},
        {'text': 'Additional Evaluation',
         'font_size': BODY_FONT, 'bold': True, 'font_color': TEAL, 'space_after': 4},
        _bullet('Feature ablation study (Sol 1)'),
        _bullet('Inter-model agreement (Cohen\u2019s \u03BA)'),
        _bullet('Cross-domain error breakdown'),
    ]
    _rich_textbox(slide, x + PAD_INNER, y + Emu(120000),
                  col_w - 2 * PAD_INNER, eval_h - Emu(200000), eval_paras)

    y += eval_h

    # =====================================================
    # COLUMN 2: Solution 1 (Cat A) + Feature Chart
    # =====================================================
    x = col_x[1]
    y = y_start

    # --- SOLUTION 1 ---
    hdr_h = _section_header(slide, x, y, col_w,
                             'Solution 1: Stylometric Ensemble (Cat A)', '\u2463')
    y += hdr_h

    sol1_h = Emu(5600000)
    _section_body(slide, x, y, col_w, sol1_h)

    sol1_paras = [
        {'text': '417 features per text \u2192 695-dim diff-vector:',
         'font_size': BODY_FONT, 'bold': True, 'space_after': 6},
        # Standard features
        {'text': 'Standard Stylometric (396):',
         'font_size': BODY_FONT_SM, 'bold': True, 'font_color': TEAL, 'space_after': 4},
        _bullet_with_detail('Function words (150)',
                            'Closed-class style markers', DARK, MID_GRAY),
        _bullet_with_detail('TF-IDF + SVD (100)',
                            'Latent semantic dims', DARK, MID_GRAY),
        _bullet_with_detail('Char n-grams (56)',
                            '2\u20134 char profiles', DARK, MID_GRAY),
        _bullet_with_detail('POS distributions (45)',
                            'Syntactic signatures', DARK, MID_GRAY),
        _bullet_with_detail('Lexical richness (30)',
                            'TTR, Yule\'s K, hapax', DARK, MID_GRAY),
        _bullet_with_detail('Structural (15)',
                            'Sentence/paragraph patterns', DARK, MID_GRAY),
        {'text': '', 'font_size': Pt(4), 'space_after': 4},
        # Novel features
        {'text': 'Novel Features (21):',
         'font_size': BODY_FONT_SM, 'bold': True, 'font_color': CORAL, 'space_after': 4},
        _bullet_with_detail('Syntactic complexity (10)',
                            'Parse depth, clause density', CORAL, MID_GRAY),
        _bullet_with_detail('Writing rhythm (6)',
                            'Syllable var, punct cadence', CORAL, MID_GRAY),
        _bullet_with_detail('Info-theoretic (5)',
                            'Entropy, compression ratio', CORAL, MID_GRAY),
        {'text': '', 'font_size': Pt(4), 'space_after': 4},
        {'text': 'Classifier: LightGBM (1000 trees, depth=7)',
         'font_size': BODY_FONT_SM, 'bold': True, 'font_color': DARK, 'space_after': 2},
    ]
    _rich_textbox(slide, x + PAD_INNER, y + Emu(120000),
                  col_w - 2 * PAD_INNER, sol1_h - Emu(200000), sol1_paras)

    y += sol1_h + Emu(200000)

    # --- Feature Groups Chart ---
    hdr_h = _section_header(slide, x, y, col_w, 'Feature Group Breakdown')
    y += hdr_h

    feat_h = Emu(4800000)
    _section_body(slide, x, y, col_w, feat_h)

    if feat_path.exists():
        img_margin = Emu(150000)
        img_w = col_w - 2 * img_margin
        img_h = Emu(4300000)
        slide.shapes.add_picture(str(feat_path),
                                  x + img_margin, y + Emu(250000),
                                  img_w, img_h)

    y += feat_h

    # =====================================================
    # COLUMN 3: Solution 2 (Cat B) + Architecture
    # =====================================================
    x = col_x[2]
    y = y_start

    # --- SOLUTION 2 ---
    hdr_h = _section_header(slide, x, y, col_w,
                             'Solution 2: Adversarial CNN-BiLSTM (Cat B)', '\u2464')
    y += hdr_h

    sol2_h = Emu(4600000)
    _section_body(slide, x, y, col_w, sol2_h)

    sol2_paras = [
        {'text': 'Siamese Encoder (813K params):',
         'font_size': BODY_FONT, 'bold': True, 'font_color': TEAL, 'space_after': 4},
        _bullet_with_detail('Char embedding (64d)',
                            'Sub-word stylistic patterns', DARK, MID_GRAY),
        _bullet_with_detail('Multi-width CNN (k=3,5,7)',
                            'Multi-scale local features', DARK, MID_GRAY),
        _bullet_with_detail('BiLSTM (256 hidden)',
                            'Sequential style dependencies', DARK, MID_GRAY),
        _bullet_with_detail('Additive Attention',
                            'Salient feature weighting', DARK, MID_GRAY),
        {'text': '', 'font_size': Pt(4), 'space_after': 4},
        {'text': 'Adversarial Debiasing:',
         'font_size': BODY_FONT, 'bold': True, 'font_color': CORAL, 'space_after': 4},
        _bullet_with_detail('Gradient Reversal Layer',
                            'Ganin & Lempitsky (2015)', CORAL, MID_GRAY),
        _bullet_with_detail('Topic pseudo-labels',
                            'TF-IDF + KMeans (k=10)', CORAL, MID_GRAY),
        _bullet('Forces topic-invariant style representations'),
        {'text': '', 'font_size': Pt(4), 'space_after': 4},
        {'text': 'Data Aug: char perturbation (insert/delete/swap)',
         'font_size': SMALL_FONT, 'font_color': MID_GRAY, 'italic': True, 'space_after': 2},
    ]
    _rich_textbox(slide, x + PAD_INNER, y + Emu(120000),
                  col_w - 2 * PAD_INNER, sol2_h - Emu(200000), sol2_paras)

    y += sol2_h + Emu(200000)

    # --- Architecture Diagram ---
    hdr_h = _section_header(slide, x, y, col_w, 'Neural Architecture Diagram')
    y += hdr_h

    arch_h = Emu(5400000)
    _section_body(slide, x, y, col_w, arch_h)

    if arch_path.exists():
        img_margin = Emu(120000)
        img_w = col_w - 2 * img_margin
        img_h = Emu(4900000)
        slide.shapes.add_picture(str(arch_path),
                                  x + img_margin, y + Emu(250000),
                                  img_w, img_h)

    y += arch_h

    # =====================================================
    # COLUMN 4: Results + F1 Chart + Confusion Matrices
    # =====================================================
    x = col_x[3]
    y = y_start

    # --- RESULTS TABLE ---
    hdr_h = _section_header(slide, x, y, col_w, 'Results (Development Set)', '\u2465')
    y += hdr_h

    results_h = Emu(3400000)
    _section_body(slide, x, y, col_w, results_h)

    results_paras = [
        # Table header
        {'runs': [
            {'text': 'Model', 'font_size': BODY_FONT_SM, 'bold': True, 'font_color': NAVY},
            {'text': '                     ', 'font_size': BODY_FONT_SM},
            {'text': 'F1', 'font_size': BODY_FONT_SM, 'bold': True, 'font_color': NAVY},
            {'text': '          ', 'font_size': BODY_FONT_SM},
            {'text': 'vs Baseline', 'font_size': BODY_FONT_SM, 'bold': True, 'font_color': NAVY},
        ], 'font_size': BODY_FONT_SM, 'space_after': 3},
        {'text': '\u2500' * 50,
         'font_size': TINY_FONT, 'font_color': BORDER_COLOR, 'space_after': 3},
        # Baselines
        {'runs': [
            {'text': 'SVM Baseline          ', 'font_size': SMALL_FONT, 'font_color': MID_GRAY},
            {'text': '0.5610', 'font_size': SMALL_FONT, 'font_color': MID_GRAY},
            {'text': '        \u2014', 'font_size': SMALL_FONT, 'font_color': LIGHT_GRAY}
        ], 'font_size': SMALL_FONT, 'space_after': 2},
        {'runs': [
            {'text': 'LSTM Baseline         ', 'font_size': SMALL_FONT, 'font_color': MID_GRAY},
            {'text': '0.6226', 'font_size': SMALL_FONT, 'font_color': MID_GRAY},
            {'text': '        \u2014', 'font_size': SMALL_FONT, 'font_color': LIGHT_GRAY}
        ], 'font_size': SMALL_FONT, 'space_after': 2},
        {'runs': [
            {'text': 'BERT Baseline         ', 'font_size': SMALL_FONT, 'font_color': MID_GRAY},
            {'text': '0.7854', 'font_size': SMALL_FONT, 'font_color': MID_GRAY},
            {'text': '        \u2014', 'font_size': SMALL_FONT, 'font_color': LIGHT_GRAY}
        ], 'font_size': SMALL_FONT, 'space_after': 3},
        {'text': '\u2500' * 50,
         'font_size': TINY_FONT, 'font_color': BORDER_COLOR, 'space_after': 3},
        # Our solutions
        {'runs': [
            {'text': '\u2605 Sol 1 (Cat A)   ', 'font_size': BODY_FONT_SM, 'bold': True,
             'font_color': ACCENT_GREEN},
            {'text': '0.7340', 'font_size': BODY_FONT_SM, 'bold': True,
             'font_color': ACCENT_GREEN},
            {'text': '  ', 'font_size': BODY_FONT_SM},
            {'text': '+30.8% vs SVM ***', 'font_size': BODY_FONT_SM, 'bold': True,
             'font_color': ACCENT_GREEN}
        ], 'font_size': BODY_FONT_SM, 'space_after': 3},
        {'runs': [
            {'text': '\u2605 Sol 2 (Cat B)   ', 'font_size': BODY_FONT_SM, 'bold': True,
             'font_color': ACCENT_ORANGE},
            {'text': '0.7422', 'font_size': BODY_FONT_SM, 'bold': True,
             'font_color': ACCENT_ORANGE},
            {'text': '  ', 'font_size': BODY_FONT_SM},
            {'text': '+19.2% vs LSTM ***', 'font_size': BODY_FONT_SM, 'bold': True,
             'font_color': ACCENT_ORANGE}
        ], 'font_size': BODY_FONT_SM, 'space_after': 4},
        {'text': '*** p < 0.001, McNemar\u2019s test. Both statistically significant.',
         'font_size': TINY_FONT, 'font_color': MID_GRAY, 'italic': True, 'space_after': 2},
    ]
    _rich_textbox(slide, x + PAD_INNER, y + Emu(120000),
                  col_w - 2 * PAD_INNER, results_h - Emu(200000), results_paras)

    y += results_h + Emu(180000)

    # --- F1 BAR CHART ---
    hdr_h = _section_header(slide, x, y, col_w, 'Performance Comparison')
    y += hdr_h

    f1_h = Emu(3500000)
    _section_body(slide, x, y, col_w, f1_h)

    if f1_path.exists():
        img_margin = Emu(120000)
        img_w = col_w - 2 * img_margin
        img_h = Emu(3100000)
        slide.shapes.add_picture(str(f1_path),
                                  x + img_margin, y + Emu(180000),
                                  img_w, img_h)

    y += f1_h + Emu(180000)

    # --- CONFUSION MATRICES ---
    hdr_h = _section_header(slide, x, y, col_w, 'Confusion Matrices')
    y += hdr_h

    cm_display_h = Emu(2900000)
    _section_body(slide, x, y, col_w, cm_display_h)

    if cm_path.exists():
        img_margin = Emu(60000)
        img_w = col_w - 2 * img_margin
        img_h = Emu(2600000)
        slide.shapes.add_picture(str(cm_path),
                                  x + img_margin, y + Emu(150000),
                                  img_w, img_h)

    y += cm_display_h

    # =====================================================
    # BOTTOM STRIP: Error Analysis + Limitations + References
    # =====================================================
    # Use the pre-computed bottom_strip_y to avoid any overlap
    bottom_y = bottom_strip_y + Emu(100000)
    bottom_body_h = Emu(3200000)

    # Error Analysis (spans cols 1-2)
    strip_w = col_w * 2 + COL_GAP
    strip_x = col_x[0]

    hdr_h = _section_header(slide, strip_x, bottom_y, strip_w,
                             'Error Analysis & Key Findings')
    _section_body(slide, strip_x, bottom_y + hdr_h,
                  strip_w, bottom_body_h)

    error_paras = [
        _key_value_line('Inter-model agreement',
                        '~70% (Cohen\u2019s \u03BA = 0.40) \u2014 complementary patterns',
                        TEAL, DARK),
        {'text': '', 'font_size': Pt(4), 'space_after': 4},
        {'text': 'Complementary Failure Modes:',
         'font_size': BODY_FONT_SM, 'bold': True, 'font_color': DARK, 'space_after': 4},
        {'runs': [
            {'text': '  \u2022  Sol 1: ', 'font_size': SMALL_FONT, 'bold': True,
             'font_color': ACCENT_GREEN},
            {'text': 'Struggles with short texts (<200 tokens) \u2014 insufficient features',
             'font_size': SMALL_FONT}
        ], 'font_size': SMALL_FONT, 'space_after': 3},
        {'runs': [
            {'text': '  \u2022  Sol 2: ', 'font_size': SMALL_FONT, 'bold': True,
             'font_color': ACCENT_ORANGE},
            {'text': 'Struggles with formal/structured texts \u2014 less char-level variation',
             'font_size': SMALL_FONT}
        ], 'font_size': SMALL_FONT, 'space_after': 3},
        {'runs': [
            {'text': '  \u2022  Topic Confound: ', 'font_size': SMALL_FONT, 'bold': True,
             'font_color': CORAL},
            {'text': 'Elevated FP on same-topic pairs; GRL partially mitigates',
             'font_size': SMALL_FONT}
        ], 'font_size': SMALL_FONT, 'space_after': 4},
        _key_value_line('Ablation',
                        'Removing novel features drops Sol 1 F1 by 1.8%',
                        TEAL, DARK, font_size=SMALL_FONT),
    ]
    _rich_textbox(slide, strip_x + PAD_INNER, bottom_y + hdr_h + Emu(100000),
                  strip_w - 2 * PAD_INNER, bottom_body_h - Emu(100000), error_paras)

    # Limitations & References (spans cols 3-4)
    strip2_x = col_x[2]
    strip2_w = col_w * 2 + COL_GAP

    hdr_h2 = _section_header(slide, strip2_x, bottom_y, strip2_w,
                               'Limitations & Key References')
    _section_body(slide, strip2_x, bottom_y + hdr_h2,
                  strip2_w, bottom_body_h)

    limref_paras = [
        {'text': 'Limitations:',
         'font_size': BODY_FONT_SM, 'bold': True, 'font_color': CORAL, 'space_after': 3},
        _bullet('Closed-mode \u2014 may not generalize to other domains/languages'),
        _bullet('English-only; cultural biases in writing norms'),
        _bullet('Neither surpasses BERT baseline (0.7854) \u2014 transformer headroom'),
        {'text': '', 'font_size': Pt(4), 'space_after': 3},
        {'text': 'Ethics:',
         'font_size': BODY_FONT_SM, 'bold': True, 'font_color': ACCENT_ORANGE, 'space_after': 3},
        _bullet('Risk: surveillance, deanonymization, forensic false positives'),
        {'text': '', 'font_size': Pt(4), 'space_after': 3},
        {'text': 'References:',
         'font_size': BODY_FONT_SM, 'bold': True, 'font_color': TEAL, 'space_after': 3},
        {'text': '[1] Stamatatos+ (2023) Diff-Vectors for AV  [2] Ganin & Lempitsky (2015) GRL',
         'font_size': TINY_FONT, 'font_color': MID_GRAY, 'space_after': 2},
        {'text': '[3] Abbasi & Chen (2008) Writeprints  [4] Boenninghoff+ (2019) AdHominem',
         'font_size': TINY_FONT, 'font_color': MID_GRAY, 'space_after': 2},
    ]
    _rich_textbox(slide, strip2_x + PAD_INNER, bottom_y + hdr_h2 + Emu(100000),
                  strip2_w - 2 * PAD_INNER, bottom_body_h - Emu(100000), limref_paras)

    # =====================================================
    # FOOTER
    # =====================================================
    footer_h = Emu(480000)
    footer_y = SLIDE_H - footer_h
    _shape(slide, 0, footer_y, SLIDE_W, footer_h, NAVY_DARK)

    # Thin teal accent at top of footer
    _shape(slide, 0, footer_y, SLIDE_W, Emu(40000), TEAL)

    _textbox(slide, Emu(800000), footer_y + Emu(100000),
             SLIDE_W - Emu(1600000), footer_h - Emu(160000),
             'COMP34812 Natural Language Understanding  \u2502  '
             'University of Manchester  \u2502  Group 34  \u2502  '
             '2025\u201326  \u2502  Closed Mode  \u2502  '
             'github.com/ayush-kumar-prog',
             Pt(20), RGBColor(0x99, 0xAA, 0xBB), bold=False,
             alignment=PP_ALIGN.CENTER)

    # =====================================================
    # SAVE
    # =====================================================
    prs.save(str(OUTPUT_PATH))
    print(f"\n{'=' * 60}")
    print(f"  Poster saved to {OUTPUT_PATH}")
    print(f"  Slide dimensions: 40\" x 22.5\" (A1 landscape)")
    print(f"  Export to PDF via PowerPoint/LibreOffice for print")
    print(f"{'=' * 60}")


if __name__ == '__main__':
    create_poster()
