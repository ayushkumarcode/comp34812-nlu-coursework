#!/usr/bin/env python3
"""
COMP34812 — A1 Landscape Academic Poster Generator (AV Track)
Group 34

Creates an A1 landscape (841mm x 594mm = 33.1" x 23.4") poster using
python-pptx with publication-quality matplotlib visualizations at 300 DPI.

Usage:
    source /tmp/poster_env/bin/activate
    python poster/create_poster_a1.py
"""

import os
import sys
from pathlib import Path

# Use poster venv if available
VENV = Path('/tmp/poster_env')
if VENV.exists():
    for p in VENV.glob('lib/*/site-packages'):
        sys.path.insert(0, str(p))

import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import matplotlib.patheffects as path_effects
import numpy as np

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ============================================================
# CONFIGURATION
# ============================================================

PROJECT_DIR = Path(__file__).parent.parent
POSTER_DIR = Path(__file__).parent
OUTPUT_PATH = POSTER_DIR / 'poster_a1.pptx'

# A1 Landscape dimensions: 841mm x 594mm = 33.1" x 23.4"
SLIDE_W = Inches(33.1)
SLIDE_H = Inches(23.4)

# Color palette — refined for print
NAVY        = RGBColor(0x1B, 0x3A, 0x5C)
NAVY_DARK   = RGBColor(0x12, 0x2A, 0x45)
TEAL        = RGBColor(0x2E, 0x86, 0xAB)
DARK_TEAL   = RGBColor(0x1E, 0x6E, 0x8E)
CORAL       = RGBColor(0xE8, 0x4D, 0x60)
LIGHT_BG    = RGBColor(0xF7, 0xF9, 0xFC)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
DARK        = RGBColor(0x2D, 0x2D, 0x2D)
MID_GRAY    = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY  = RGBColor(0xBB, 0xBB, 0xBB)
ACCENT_GREEN  = RGBColor(0x27, 0xAE, 0x60)
ACCENT_ORANGE = RGBColor(0xE6, 0x7E, 0x22)
TABLE_ALT_BG  = RGBColor(0xEB, 0xF0, 0xF5)
BORDER_COLOR  = RGBColor(0xDD, 0xDD, 0xDD)
GOLD          = RGBColor(0xF3, 0x9C, 0x12)
SOFT_GREEN_BG = RGBColor(0xE8, 0xF5, 0xE9)
SOFT_ORANGE_BG = RGBColor(0xFF, 0xF3, 0xE0)
KEY_FINDING_BG = RGBColor(0xFF, 0xF8, 0xE1)

# Hex strings for matplotlib
C_NAVY = '#1B3A5C'
C_TEAL = '#2E86AB'
C_CORAL = '#E84D60'
C_GREEN = '#27AE60'
C_ORANGE = '#E67E22'
C_DARK = '#2D2D2D'
C_GRAY = '#95A5A6'
C_LIGHTBG = '#F7F9FC'

# Font sizes (calibrated for A1 print — readable from 1.5m)
TITLE_FONT     = Pt(60)
SUBTITLE_FONT  = Pt(28)
HEADER_FONT    = Pt(30)
BODY_FONT      = Pt(22)
BODY_FONT_SM   = Pt(20)
SMALL_FONT     = Pt(18)
TINY_FONT      = Pt(16)
METRIC_FONT    = Pt(24)

# Layout constants
MARGIN     = Emu(450000)
COL_GAP    = Emu(320000)
SEC_GAP    = Emu(180000)
PAD_INNER  = Emu(160000)


# ============================================================
# CHART GENERATION
# ============================================================

def load_predictions():
    """Load ground truth and best predictions for AV task."""
    gt = []
    gt_path = PROJECT_DIR / 'baseline_extracted' / 'nlu_bundle-feature-unified-local-scorer' / \
              'local_scorer' / 'reference_data' / 'NLU_SharedTask_AV_dev.solution'
    with open(gt_path) as f:
        for line in f:
            gt.append(int(line.strip()))

    def load_pred(path):
        preds = []
        with open(path) as f:
            for line in f:
                preds.append(int(line.strip()))
        return preds

    pred_dir = PROJECT_DIR / 'predictions'
    pred_a = load_pred(pred_dir / 'av_Group_34_A_lgbm.csv')
    pred_b = load_pred(pred_dir / 'av_Group_34_B_v3.csv')

    if len(pred_a) > len(gt):
        pred_a = pred_a[1:]
    if len(pred_b) > len(gt):
        pred_b = pred_b[1:]

    return np.array(gt), np.array(pred_a), np.array(pred_b)


def _set_chart_style():
    """Set consistent matplotlib style for all charts."""
    plt.rcParams.update({
        'font.family': 'sans-serif',
        'font.sans-serif': ['Helvetica Neue', 'Helvetica', 'Arial', 'DejaVu Sans'],
        'font.size': 14,
        'axes.labelsize': 18,
        'axes.titlesize': 22,
        'xtick.labelsize': 15,
        'ytick.labelsize': 15,
        'figure.facecolor': 'white',
        'axes.facecolor': 'white',
        'axes.edgecolor': '#CCCCCC',
        'axes.grid': False,
        'grid.alpha': 0.3,
        'grid.linestyle': '--',
        'legend.framealpha': 0.95,
        'legend.edgecolor': '#CCCCCC',
    })


def generate_f1_chart():
    """Generate polished F1 comparison bar chart — clean, no clutter."""
    _set_chart_style()

    models = ['SVM\n(Baseline)', 'LSTM\n(Baseline)', 'BERT\n(Baseline)',
              'Sol 1\n(Cat A)', 'Sol 2\n(Cat B)']
    scores = [0.5610, 0.6226, 0.7854, 0.7340, 0.7422]

    fig, ax = plt.subplots(figsize=(11, 7.5))

    # Colors: gray for baselines, distinct for our solutions
    bar_colors = ['#B0BEC5', '#B0BEC5', '#B0BEC5', '#27AE60', '#E67E22']
    edge_colors = ['#90A4AE', '#90A4AE', '#90A4AE', '#1B8C4F', '#C56E1A']

    bars = ax.bar(models, scores, color=bar_colors, edgecolor=edge_colors,
                  linewidth=2, width=0.58, zorder=3)

    # Value labels on top of each bar
    for bar, score in zip(bars, scores):
        txt = ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 0.008,
                      f'{score:.4f}', ha='center', va='bottom',
                      fontsize=19, fontweight='bold', color=C_DARK)
        txt.set_path_effects([
            path_effects.withStroke(linewidth=3, foreground='white')
        ])

    # Improvement annotations — inside bars at bottom
    ax.text(3, 0.08, '+30.8%\nvs SVM ***', fontsize=12,
            fontweight='bold', color='white', ha='center', va='bottom',
            bbox=dict(boxstyle='round,pad=0.25', facecolor=C_GREEN,
                      edgecolor='white', alpha=0.85, linewidth=1.5))

    ax.text(4, 0.08, '+19.2%\nvs LSTM ***', fontsize=12,
            fontweight='bold', color='white', ha='center', va='bottom',
            bbox=dict(boxstyle='round,pad=0.25', facecolor=C_ORANGE,
                      edgecolor='white', alpha=0.85, linewidth=1.5))

    # BERT baseline reference line (no label -- already has score above bar)
    ax.axhline(y=0.7854, color='#BDC3C7', linestyle=':', linewidth=1.5, alpha=0.4)

    ax.set_ylabel('Macro F1 Score', fontsize=20, fontweight='bold',
                   labelpad=12, color=C_DARK)
    ax.set_ylim(0, 0.88)
    ax.set_xlim(-0.5, 4.5)

    ax.yaxis.grid(True, linestyle='--', alpha=0.3, zorder=0, color='#DDD')
    ax.set_axisbelow(True)

    ax.tick_params(axis='both', labelsize=15, colors=C_DARK)
    ax.tick_params(axis='x', length=0, pad=10)

    for spine in ['top', 'right']:
        ax.spines[spine].set_visible(False)
    for spine in ['left', 'bottom']:
        ax.spines[spine].set_color('#CCCCCC')
        ax.spines[spine].set_linewidth(1.5)

    # Legend
    baseline_patch = mpatches.Patch(facecolor='#B0BEC5', edgecolor='#90A4AE',
                                     linewidth=1.5, label='Baselines')
    sol1_patch = mpatches.Patch(facecolor=C_GREEN, edgecolor='#1B8C4F',
                                 linewidth=1.5, label='Sol 1 (Cat A)')
    sol2_patch = mpatches.Patch(facecolor=C_ORANGE, edgecolor='#C56E1A',
                                 linewidth=1.5, label='Sol 2 (Cat B)')
    ax.legend(handles=[baseline_patch, sol1_patch, sol2_patch],
              fontsize=14, loc='upper left', framealpha=0.95,
              edgecolor='#CCCCCC', fancybox=True)

    ax.text(0.98, 0.97, '*** p < 0.001 (McNemar\'s test)',
            transform=ax.transAxes, ha='right', va='top', fontsize=12,
            color='#888888', style='italic')

    plt.tight_layout(pad=1.2)
    path = POSTER_DIR / 'f1_chart_a1.png'
    plt.savefig(str(path), dpi=300, bbox_inches='tight',
                facecolor='white', edgecolor='none')
    plt.close()
    print(f"  [OK] F1 chart -> {path}")
    return path


def generate_confusion_matrices(gt, pred_a, pred_b):
    """Generate side-by-side confusion matrices with real data."""
    _set_chart_style()
    from sklearn.metrics import confusion_matrix, precision_score, recall_score, f1_score

    cm_a = confusion_matrix(gt, pred_a)
    cm_b = confusion_matrix(gt, pred_b)

    fig, axes = plt.subplots(1, 2, figsize=(15, 6.5))

    titles = ['Sol 1 (Cat A) \u2014 Stylometric LightGBM',
              'Sol 2 (Cat B) \u2014 Adversarial CNN-BiLSTM']
    cms = [cm_a, cm_b]
    cmaps = ['Blues', 'Oranges']
    preds_list = [pred_a, pred_b]

    for ax, cm, title, cmap, preds in zip(axes, cms, titles, cmaps, preds_list):
        cm_norm = cm.astype(float) / cm.sum(axis=1, keepdims=True)
        im = ax.imshow(cm_norm, cmap=cmap, interpolation='nearest',
                        aspect='equal', vmin=0, vmax=1)

        ax.set_xticks([0, 1])
        ax.set_yticks([0, 1])
        ax.set_xticklabels(['Diff.\nAuthor', 'Same\nAuthor'], fontsize=14)
        ax.set_yticklabels(['Diff.\nAuthor', 'Same\nAuthor'], fontsize=14)
        ax.set_xlabel('Predicted', fontsize=17, fontweight='bold', labelpad=10)
        ax.set_ylabel('True', fontsize=17, fontweight='bold', labelpad=10)
        ax.set_title(title, fontsize=17, fontweight='bold', pad=14, color=C_NAVY)

        for i in range(2):
            for j in range(2):
                val = cm[i, j]
                pct = cm_norm[i, j] * 100
                color = 'white' if cm_norm[i, j] > 0.5 else C_DARK
                ax.text(j, i, f'{val:,}\n({pct:.1f}%)', ha='center', va='center',
                        fontsize=18, fontweight='bold', color=color)

        # Highlight diagonal
        for k in range(2):
            accent = '#27AE60' if cmap == 'Blues' else '#E67E22'
            rect = plt.Rectangle((k-0.5, k-0.5), 1, 1, linewidth=2.5,
                                  edgecolor=accent, facecolor='none',
                                  linestyle='--', alpha=0.6)
            ax.add_patch(rect)

        prec = precision_score(gt, preds, average='macro')
        rec = recall_score(gt, preds, average='macro')
        f1 = f1_score(gt, preds, average='macro')
        acc = np.mean(gt == preds)
        ax.text(0.5, -0.22,
                f'Acc={acc:.3f}  P={prec:.3f}  R={rec:.3f}  F1={f1:.4f}',
                transform=ax.transAxes, ha='center', fontsize=13,
                color='#444', fontweight='bold',
                bbox=dict(boxstyle='round,pad=0.35', facecolor='#F5F5F5',
                          edgecolor='#CCC', alpha=0.9))

    plt.tight_layout(w_pad=3.5, pad=2)
    path = POSTER_DIR / 'cm_a1.png'
    plt.savefig(str(path), dpi=300, bbox_inches='tight',
                facecolor='white', edgecolor='none')
    plt.close()
    print(f"  [OK] Confusion matrices -> {path}")
    return path


def generate_architecture_diagram():
    """Generate ACCURATE architecture diagram for Sol 2.

    Verified architecture from submission/src/models/av_cat_b_model.py:
    - vocab_size=97, char_emb_dim=32
    - Multi-width CNN: k=3,5,7, each 128 filters -> 384 total
    - MaxPool1d(k=3, stride=3)
    - BiLSTM: hidden=128, output=256 (bidirectional)
    - Additive (Bahdanau) Attention -> 256d
    - Projection: 256->128d with ReLU+Dropout(0.3)
    - Comparison: [v1, v2, |v1-v2|, v1*v2] = 512d
    - MLP: 512->256->64->1
    - GRL: topic head on individual embeddings (proj_dim->64->10)
    """
    _set_chart_style()
    plt.rcParams['font.sans-serif'] = ['DejaVu Sans', 'Arial']

    fig, ax = plt.subplots(figsize=(13, 10))
    ax.set_xlim(-0.5, 10.5)
    ax.set_ylim(-0.5, 11.5)
    ax.axis('off')
    fig.patch.set_facecolor('white')

    def draw_box(x, y, w, h, text, color, text_color='white',
                 fontsize=13, alpha=1.0, style='round,pad=0.15',
                 linewidth=2.5, edgecolor=None):
        ec = edgecolor or color
        # Shadow
        shadow = mpatches.FancyBboxPatch(
            (x - w/2 + 0.04, y - h/2 - 0.04), w, h,
            boxstyle=style, facecolor='#00000008',
            edgecolor='none', linewidth=0)
        ax.add_patch(shadow)
        rect = mpatches.FancyBboxPatch(
            (x - w/2, y - h/2), w, h,
            boxstyle=style, facecolor=color,
            edgecolor=ec, linewidth=linewidth, alpha=alpha)
        ax.add_patch(rect)
        ax.text(x, y, text, ha='center', va='center',
                fontsize=fontsize, fontweight='bold', color=text_color,
                zorder=10)

    def arrow(x1, y1, x2, y2, color='#888', width=1.8, style='->'):
        ax.annotate('', xy=(x2, y2), xytext=(x1, y1),
                    arrowprops=dict(arrowstyle=style, color=color,
                                    lw=width, connectionstyle='arc3,rad=0'))

    # Title
    ax.text(5, 11.1, 'Sol 2: Adversarial Style-Content Disentanglement',
            ha='center', fontsize=21, fontweight='bold', color=C_NAVY)
    ax.text(5, 10.75, 'Siamese Char-CNN + BiLSTM + Attention + GRL (813,931 params)',
            ha='center', fontsize=13, color='#777', style='italic')

    # Input layer
    draw_box(2.5, 10.0, 3.0, 0.55, 'Text 1 (max 1500 chars)', '#E3F2FD', C_NAVY, 13,
             edgecolor='#90CAF9')
    draw_box(7.5, 10.0, 3.0, 0.55, 'Text 2 (max 1500 chars)', '#E3F2FD', C_NAVY, 13,
             edgecolor='#90CAF9')

    # Char Embedding (CORRECT: 32d)
    draw_box(2.5, 8.95, 3.0, 0.55, 'Char Embed (32d)', '#42A5F5', 'white', 13)
    draw_box(7.5, 8.95, 3.0, 0.55, 'Char Embed (32d)', '#42A5F5', 'white', 13)

    # Multi-width CNN
    for offset, ks in [(-0.85, '3'), (0, '5'), (0.85, '7')]:
        draw_box(2.5 + offset, 7.75, 0.7, 0.5, f'k={ks}', '#1E88E5', 'white', 11,
                 style='round,pad=0.08', linewidth=1.5)
        draw_box(7.5 + offset, 7.75, 0.7, 0.5, f'k={ks}', '#1E88E5', 'white', 11,
                 style='round,pad=0.08', linewidth=1.5)

    ax.text(2.5, 8.25, 'Multi-Width CNN (128 each)', ha='center', fontsize=11,
            fontweight='bold', color=C_NAVY)
    ax.text(7.5, 8.25, 'Multi-Width CNN (128 each)', ha='center', fontsize=11,
            fontweight='bold', color=C_NAVY)

    # MaxPool + concat annotation
    ax.text(2.5, 7.3, 'MaxPool(3) + Concat = 384d', ha='center', fontsize=11,
            color='#666', style='italic', fontweight='bold')
    ax.text(7.5, 7.3, 'MaxPool(3) + Concat = 384d', ha='center', fontsize=11,
            color='#666', style='italic', fontweight='bold')

    # BiLSTM (CORRECT: 128 hidden, 256 output)
    draw_box(2.5, 6.55, 3.2, 0.55, 'BiLSTM (128h -> 256d)', '#1565C0', 'white', 12)
    draw_box(7.5, 6.55, 3.2, 0.55, 'BiLSTM (128h -> 256d)', '#1565C0', 'white', 12)

    # Attention
    draw_box(2.5, 5.5, 3.0, 0.55, 'Additive Attention', '#0D47A1', 'white', 13)
    draw_box(7.5, 5.5, 3.0, 0.55, 'Additive Attention', '#0D47A1', 'white', 13)

    # Projection
    draw_box(2.5, 4.5, 2.5, 0.5, 'Projection (128d)', '#5C6BC0', 'white', 12)
    draw_box(7.5, 4.5, 2.5, 0.5, 'Projection (128d)', '#5C6BC0', 'white', 12)

    # v1, v2 labels
    ax.text(2.5, 3.95, 'v\u2081', ha='center', fontsize=16,
            fontweight='bold', color=C_NAVY, style='italic')
    ax.text(7.5, 3.95, 'v\u2082', ha='center', fontsize=16,
            fontweight='bold', color=C_NAVY, style='italic')

    # Shared weights indicator
    ax.annotate('', xy=(4.0, 8.95), xytext=(6.0, 8.95),
                arrowprops=dict(arrowstyle='<->', color=C_CORAL, lw=2.5,
                                connectionstyle='arc3,rad=0'))
    ax.text(5, 9.2, 'Shared\nWeights', ha='center', fontsize=10,
            color=C_CORAL, fontweight='bold',
            bbox=dict(boxstyle='round,pad=0.15', facecolor='white',
                      edgecolor=C_CORAL, alpha=0.9))

    # Arrows through encoder
    for y_top, y_bot in [(9.72, 9.22), (8.67, 8.25), (7.48, 6.82),
                          (6.27, 5.77), (5.22, 4.75)]:
        arrow(2.5, y_top, 2.5, y_bot, '#AAA', 1.3)
        arrow(7.5, y_top, 7.5, y_bot, '#AAA', 1.3)

    # Merge/comparison
    draw_box(5, 3.15, 4.8, 0.65, '[v\u2081, v\u2082, |v\u2081\u2212v\u2082|, v\u2081\u2299v\u2082] = 512d',
             '#7B1FA2', 'white', 14)

    arrow(2.5, 3.85, 3.5, 3.47, '#888', 1.5)
    arrow(7.5, 3.85, 6.5, 3.47, '#888', 1.5)

    # Two branches
    draw_box(3.0, 1.65, 3.6, 0.6, 'MLP (256->64->1)', C_GREEN, 'white', 13, linewidth=2.5)
    ax.text(3.0, 1.0, 'Same / Different Author',
            ha='center', fontsize=11, color=C_GREEN, fontweight='bold')

    draw_box(7.2, 1.65, 3.6, 0.6, 'GRL -> Topic Head', C_CORAL, 'white', 13, linewidth=2.5)
    ax.text(7.2, 1.0, 'Gradient Reversal Layer',
            ha='center', fontsize=11, color=C_CORAL, fontweight='bold')
    ax.text(7.2, 0.6, 'lambda ramp 0->0.05 over 20 epochs',
            ha='center', fontsize=9, color='#999', style='italic')

    arrow(4.0, 2.82, 3.0, 1.95, C_GREEN, 2)
    arrow(6.0, 2.82, 7.2, 1.95, C_CORAL, 2)

    # Dimension annotations on left side
    dims = [(8.95, '32'), (7.75, '384'), (6.55, '256'), (5.5, '256'), (4.5, '128')]
    for y_pos, dim_text in dims:
        ax.text(0.5, y_pos, dim_text, ha='center', fontsize=9, color='#AAA',
                bbox=dict(boxstyle='round', facecolor='#F8F8F8', edgecolor='#DDD'))

    plt.tight_layout(pad=0.8)
    path = POSTER_DIR / 'arch_diagram_a1.png'
    plt.savefig(str(path), dpi=300, bbox_inches='tight',
                facecolor='white', edgecolor='none')
    plt.close()
    print(f"  [OK] Architecture diagram -> {path}")
    return path


def generate_feature_groups_chart():
    """Generate horizontal bar chart of feature groups.

    Per-text features: 417 (summed from 9 groups below).
    After diff-vector |f(t1) - f(t2)| and concat: 695 dimensions.
    But the chart shows the 9 feature GROUPS and their per-text counts.
    """
    _set_chart_style()

    groups = [
        'Function Words', 'TF-IDF + SVD', 'Character N-gram',
        'POS Tags', 'Lexical Richness', 'Structural',
        'Syntactic Complexity', 'Writing Rhythm', 'Info-Theoretic'
    ]
    counts = [150, 100, 56, 45, 30, 15, 10, 6, 5]
    is_novel = [False, False, False, False, False, False, True, True, True]
    colors = [C_CORAL if n else C_TEAL for n in is_novel]
    edge_colors = ['#C43350' if n else '#1E6E8E' for n in is_novel]

    fig, ax = plt.subplots(figsize=(10, 6.5))
    y_pos = np.arange(len(groups))

    bars = ax.barh(y_pos, counts, color=colors, edgecolor=edge_colors,
                   linewidth=1.5, height=0.55, zorder=3)

    for bar, count, novel in zip(bars, counts, is_novel):
        label_color = C_CORAL if novel else C_TEAL
        ax.text(bar.get_width() + 2.5, bar.get_y() + bar.get_height()/2,
                f'{count}', va='center', fontsize=16, fontweight='bold',
                color=label_color)

    ax.set_yticks(y_pos)
    ax.set_yticklabels(groups, fontsize=15, fontweight='bold')
    ax.set_xlabel('Features per Group', fontsize=18, fontweight='bold',
                  labelpad=10, color=C_DARK)
    ax.set_title('9 Feature Groups (417 per text -> 695d diff-vector)',
                 fontsize=20, fontweight='bold', pad=16, color=C_NAVY)
    ax.set_xlim(0, 185)
    ax.invert_yaxis()

    ax.xaxis.grid(True, linestyle='--', alpha=0.3, zorder=0, color='#DDD')
    for spine in ['top', 'right']:
        ax.spines[spine].set_visible(False)
    ax.spines['left'].set_color('#CCCCCC')
    ax.spines['bottom'].set_color('#CCCCCC')

    std_patch = mpatches.Patch(facecolor=C_TEAL, edgecolor='#1E6E8E',
                                linewidth=1.5, label='Standard Stylometric (396)')
    novel_patch = mpatches.Patch(facecolor=C_CORAL, edgecolor='#C43350',
                                  linewidth=1.5, label='Novel Features (21)')
    ax.legend(handles=[std_patch, novel_patch], fontsize=14,
              loc='lower right', framealpha=0.95, edgecolor='#CCCCCC',
              fancybox=True)

    plt.tight_layout(pad=1.2)
    path = POSTER_DIR / 'feature_groups_a1.png'
    plt.savefig(str(path), dpi=300, bbox_inches='tight',
                facecolor='white', edgecolor='none')
    plt.close()
    print(f"  [OK] Feature groups chart -> {path}")
    return path


# ============================================================
# POSTER LAYOUT HELPERS
# ============================================================

def _shape(slide, left, top, width, height, fill_color,
           line_color=None, line_width=None):
    """Add a filled rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def _rounded_rect(slide, left, top, width, height, fill_color,
                  line_color=None, line_width=None):
    """Add a rounded rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    shape.adjustments[0] = 0.04
    return shape


def _textbox(slide, left, top, width, height, text, font_size,
             font_color=DARK, bold=False, alignment=PP_ALIGN.LEFT,
             font_name='Calibri', line_spacing=1.15, italic=False):
    """Add a simple text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    p.font.italic = italic
    p.alignment = alignment
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    if line_spacing:
        p.line_spacing = line_spacing
    return txBox


def _rich_textbox(slide, left, top, width, height, paragraphs_data,
                  font_name='Calibri'):
    """Add a text box with multiple styled paragraphs."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, pdata in enumerate(paragraphs_data):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        runs = pdata.get('runs', None)
        if runs:
            for j, run_data in enumerate(runs):
                if j == 0 and p.runs:
                    run = p.runs[0]
                else:
                    run = p.add_run()
                run.text = run_data.get('text', '')
                run.font.size = run_data.get('font_size', pdata.get('font_size', BODY_FONT))
                run.font.color.rgb = run_data.get('font_color', pdata.get('font_color', DARK))
                run.font.bold = run_data.get('bold', pdata.get('bold', False))
                run.font.name = run_data.get('font_name', font_name)
                if run_data.get('italic', False):
                    run.font.italic = True
        else:
            p.text = pdata.get('text', '')
            p.font.size = pdata.get('font_size', BODY_FONT)
            p.font.color.rgb = pdata.get('font_color', DARK)
            p.font.bold = pdata.get('bold', False)
            p.font.name = pdata.get('font_name', font_name)
            if pdata.get('italic', False):
                p.font.italic = True

        p.alignment = pdata.get('alignment', PP_ALIGN.LEFT)
        p.space_after = Pt(pdata.get('space_after', 4))
        p.space_before = Pt(pdata.get('space_before', 0))
        ls = pdata.get('line_spacing', 1.15)
        if ls:
            p.line_spacing = ls

    return txBox


def _section_header(slide, left, top, width, title, number=None):
    """Add a clean section header with teal bar and number."""
    h = Emu(520000)

    # Main teal bar
    _shape(slide, left, top, width, h, TEAL)

    # Dark accent on left
    accent_w = Emu(60000)
    _shape(slide, left, top, accent_w, h, DARK_TEAL)

    # Optional number circle
    if number:
        _textbox(slide, left + Emu(100000), top + Emu(60000),
                 Emu(320000), h - Emu(120000),
                 number, Pt(24), WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)
        text_left = left + Emu(360000)
    else:
        text_left = left + Emu(180000)

    _textbox(slide, text_left, top + Emu(65000),
             width - Emu(250000), h - Emu(130000),
             title, HEADER_FONT, WHITE, bold=True)

    # Bottom accent line
    _shape(slide, left, top + h - Emu(35000), width, Emu(35000), DARK_TEAL)

    return h


def _section_body(slide, left, top, width, height):
    """Add section body background."""
    return _rounded_rect(slide, left, top, width, height,
                          LIGHT_BG, BORDER_COLOR, Pt(1.5))


def _kv(label, value, label_color=DARK, value_color=DARK,
        font_size=BODY_FONT, bold_label=True):
    return {
        'runs': [
            {'text': f'{label}: ', 'font_size': font_size, 'bold': bold_label,
             'font_color': label_color},
            {'text': value, 'font_size': font_size, 'font_color': value_color}
        ],
        'font_size': font_size,
        'space_after': 5
    }


def _bullet(text, font_size=BODY_FONT_SM, color=DARK, bold=False):
    return {
        'runs': [
            {'text': '  \u2022  ', 'font_size': font_size, 'font_color': MID_GRAY},
            {'text': text, 'font_size': font_size, 'font_color': color, 'bold': bold}
        ],
        'font_size': font_size,
        'space_after': 3
    }


def _bullet_detail(label, detail, label_color=DARK, detail_color=MID_GRAY,
                   font_size=BODY_FONT_SM):
    return {
        'runs': [
            {'text': '  \u2022  ', 'font_size': font_size, 'font_color': MID_GRAY},
            {'text': label, 'font_size': font_size, 'font_color': label_color, 'bold': True},
            {'text': f' \u2014 {detail}', 'font_size': Pt(18), 'font_color': detail_color}
        ],
        'font_size': font_size,
        'space_after': 3
    }


# ============================================================
# MAIN POSTER CREATION — A1 LANDSCAPE (3 columns)
# ============================================================

def create_poster():
    """Create A1 landscape AV track academic poster."""
    print("=" * 60)
    print("  Creating AV Track Academic Poster (A1 Landscape)")
    print("  Dimensions: 33.1\" x 23.4\" (841mm x 594mm)")
    print("=" * 60)

    # Load real data
    print("\nLoading prediction data...")
    gt, pred_a, pred_b = load_predictions()
    print(f"  Ground truth: {len(gt)} samples")
    print(f"  Pred A (Cat A): {len(pred_a)} samples")
    print(f"  Pred B (Cat B): {len(pred_b)} samples")

    # Generate all visualizations
    print("\nGenerating visualizations at 300 DPI...")
    f1_path = generate_f1_chart()
    cm_path = generate_confusion_matrices(gt, pred_a, pred_b)
    arch_path = generate_architecture_diagram()
    feat_path = generate_feature_groups_chart()

    # Create presentation
    print("\nBuilding A1 poster layout (3 columns)...")
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0xEE, 0xF1, 0xF5)

    # =====================================================
    # TITLE BANNER
    # =====================================================
    title_h = Emu(2100000)  # ~2.3"

    _shape(slide, 0, 0, SLIDE_W, title_h, NAVY)
    _shape(slide, 0, 0, SLIDE_W, Emu(50000), NAVY_DARK)
    _shape(slide, 0, title_h - Emu(65000), SLIDE_W, Emu(65000), TEAL)
    _shape(slide, 0, title_h - Emu(75000), SLIDE_W, Emu(10000), CORAL)

    # Title
    _textbox(slide, Emu(500000), Emu(120000),
             SLIDE_W - Emu(1000000), Emu(1100000),
             'Authorship Verification: Combining Stylometric Feature\n'
             'Engineering with Adversarial Neural Style-Content Disentanglement',
             Pt(54), WHITE, bold=True, alignment=PP_ALIGN.CENTER,
             line_spacing=1.05)

    # Subtitle
    _textbox(slide, Emu(500000), Emu(1350000),
             SLIDE_W - Emu(1000000), Emu(320000),
             'COMP34812 Natural Language Understanding  |  Group 34  |  '
             'University of Manchester  |  2025\u201326',
             Pt(26), RGBColor(0xBB, 0xCC, 0xDD), bold=False,
             alignment=PP_ALIGN.CENTER)

    # Key result highlight in title bar
    _textbox(slide, Emu(500000), Emu(1700000),
             SLIDE_W - Emu(1000000), Emu(250000),
             'Cat A: F1=0.7340 (+30.8% vs SVM, p<0.001)   |   '
             'Cat B: F1=0.7422 (+19.2% vs LSTM, p<0.001)',
             Pt(22), RGBColor(0xFF, 0xD5, 0x4F), bold=True,
             alignment=PP_ALIGN.CENTER)

    # =====================================================
    # 3-COLUMN LAYOUT
    # =====================================================
    usable_w = SLIDE_W - 2 * MARGIN
    n_cols = 3
    col_w = (usable_w - (n_cols - 1) * COL_GAP) // n_cols
    y_start = title_h + Emu(130000)

    footer_h = Emu(400000)
    bottom_limit = SLIDE_H - footer_h - Emu(80000)

    col_x = []
    for i in range(n_cols):
        col_x.append(MARGIN + i * (col_w + COL_GAP))

    # Column dividers
    for i in range(1, n_cols):
        div_x = col_x[i] - COL_GAP // 2
        _shape(slide, div_x, y_start, Emu(10000),
               bottom_limit - y_start,
               RGBColor(0xDD, 0xDD, 0xDD))

    # =====================================================
    # COLUMN 1: Task + Dataset + Solution 1 + Feature Chart
    # =====================================================
    x = col_x[0]
    y = y_start

    # --- TASK & MOTIVATION ---
    hdr_h = _section_header(slide, x, y, col_w, 'Task & Motivation', '1')
    y += hdr_h

    intro_h = Emu(2600000)
    _section_body(slide, x, y, col_w, intro_h)

    intro_paras = [
        {'text': 'Authorship Verification (AV): given two texts, determine if '
                 'they share the same author. Binary classification task '
                 'requiring style-sensitive, content-invariant models.',
         'font_size': BODY_FONT, 'space_after': 8},
        {'text': 'Core Challenge: Style-Content Confound',
         'font_size': BODY_FONT, 'bold': True, 'font_color': CORAL, 'space_after': 5},
        {'text': 'Models may learn topic similarity as a proxy for shared '
                 'authorship, producing brittle correlations that fail to generalize.',
         'font_size': BODY_FONT, 'space_after': 8},
        {'text': 'Our Two Solutions:',
         'font_size': BODY_FONT, 'bold': True, 'font_color': TEAL, 'space_after': 5},
        _bullet_detail('Sol 1 (Cat A)', 'Stylometric LightGBM with 695-dim diff-vectors',
                        ACCENT_GREEN),
        _bullet_detail('Sol 2 (Cat B)', 'Siamese CNN-BiLSTM with adversarial debiasing (GRL)',
                        ACCENT_ORANGE),
    ]
    _rich_textbox(slide, x + PAD_INNER, y + Emu(110000),
                  col_w - 2 * PAD_INNER, intro_h - Emu(180000), intro_paras)
    y += intro_h + SEC_GAP

    # --- DATASET ---
    hdr_h = _section_header(slide, x, y, col_w, 'Dataset Summary', '2')
    y += hdr_h

    dataset_h = Emu(2050000)
    _section_body(slide, x, y, col_w, dataset_h)

    dataset_paras = [
        _kv('Training pairs', '27,643 (51% same-author, 49% different)', TEAL),
        _kv('Dev pairs', '5,993 (same distribution)', TEAL),
        _kv('Domains', 'Cross-domain (emails, blogs, reviews, essays)', TEAL),
        _kv('Mode', 'Closed \u2014 only provided training data', TEAL),
        {'text': '', 'font_size': Pt(4), 'space_after': 4},
        _bullet('Near-balanced: macro F1 closely tracks accuracy'),
        _bullet('Each pair: 2 texts (variable length, 50-10,000 tokens)'),
    ]
    _rich_textbox(slide, x + PAD_INNER, y + Emu(110000),
                  col_w - 2 * PAD_INNER, dataset_h - Emu(180000), dataset_paras)
    y += dataset_h + SEC_GAP

    # --- SOLUTION 1 ---
    hdr_h = _section_header(slide, x, y, col_w,
                             'Sol 1: Stylometric Ensemble (Cat A)', '3')
    y += hdr_h

    sol1_h = Emu(4600000)
    _section_body(slide, x, y, col_w, sol1_h)

    sol1_paras = [
        {'text': '417 features per text across 9 groups:',
         'font_size': BODY_FONT, 'bold': True, 'space_after': 7},
        {'text': 'Standard Stylometric (396):',
         'font_size': BODY_FONT, 'bold': True, 'font_color': TEAL, 'space_after': 4},
        _bullet_detail('Function words (150)',
                        'Closed-class words as style markers'),
        _bullet_detail('TF-IDF + SVD (100)',
                        'Latent semantic dimensions'),
        _bullet_detail('Character n-grams (56)',
                        '2\u20134 char frequency profiles'),
        _bullet_detail('POS tags (45)',
                        'Syntactic preference signatures'),
        _bullet_detail('Lexical richness (30)',
                        'TTR, Yule\'s K, hapax ratios'),
        _bullet_detail('Structural (15)',
                        'Sentence/paragraph length patterns'),
        {'text': '', 'font_size': Pt(3), 'space_after': 4},
        {'text': 'Novel Features (21):',
         'font_size': BODY_FONT, 'bold': True, 'font_color': CORAL, 'space_after': 4},
        _bullet_detail('Syntactic complexity (10)',
                        'Parse tree depth, clause density', CORAL, MID_GRAY),
        _bullet_detail('Writing rhythm (6)',
                        'Syllable variance, punctuation cadence', CORAL, MID_GRAY),
        _bullet_detail('Info-theoretic (5)',
                        'Entropy rate, compression ratio', CORAL, MID_GRAY),
        {'text': '', 'font_size': Pt(3), 'space_after': 5},
        {'text': 'Representation: diff-vector |f(t1) - f(t2)| -> 695 dims',
         'font_size': BODY_FONT, 'bold': True, 'font_color': DARK, 'space_after': 3},
        {'text': 'Classifier: LightGBM (1000 trees, max_depth=7)',
         'font_size': BODY_FONT, 'bold': True, 'font_color': DARK, 'space_after': 3},
    ]
    _rich_textbox(slide, x + PAD_INNER, y + Emu(110000),
                  col_w - 2 * PAD_INNER, sol1_h - Emu(180000), sol1_paras)
    y += sol1_h + SEC_GAP

    # --- FEATURE GROUPS CHART ---
    hdr_h = _section_header(slide, x, y, col_w, 'Feature Group Breakdown')
    y += hdr_h

    feat_chart_h = Emu(4200000)
    _section_body(slide, x, y, col_w, feat_chart_h)

    if feat_path.exists():
        img_margin = Emu(120000)
        img_w = col_w - 2 * img_margin
        img_h = Emu(3800000)
        slide.shapes.add_picture(str(feat_path),
                                  x + img_margin, y + Emu(180000),
                                  img_w, img_h)
    y += feat_chart_h + SEC_GAP

    # --- KEY INSIGHT CALLOUT (uses remaining column 1 space) ---
    remaining_c1 = bottom_limit - y
    if remaining_c1 > Emu(600000):
        insight_h = min(remaining_c1, Emu(2200000))
        _rounded_rect(slide, x, y, col_w, insight_h,
                      KEY_FINDING_BG, GOLD, Pt(2))

        insight_paras = [
            {'text': 'Key Insight',
             'font_size': Pt(24), 'bold': True, 'font_color': NAVY, 'space_after': 6},
            {'text': 'The diff-vector |f(t\u2081) \u2212 f(t\u2082)| representation is '
                     'central to Cat A: each feature captures authorial style per text, '
                     'and the absolute difference measures stylistic distance. '
                     'This makes the representation symmetric and '
                     'content-agnostic by construction.',
             'font_size': BODY_FONT_SM, 'space_after': 6},
            {'text': 'The 21 novel features (syntactic complexity, writing rhythm, '
                     'information-theoretic measures) contribute a statistically '
                     'significant +1.8% F1 gain over standard stylometric features alone.',
             'font_size': BODY_FONT_SM, 'font_color': CORAL, 'bold': True, 'space_after': 3},
        ]
        _rich_textbox(slide, x + PAD_INNER, y + Emu(100000),
                      col_w - 2 * PAD_INNER, insight_h - Emu(160000), insight_paras)
        y += insight_h

    # =====================================================
    # COLUMN 2: Solution 2 + Architecture + Results
    # =====================================================
    x = col_x[1]
    y = y_start

    # --- SOLUTION 2 ---
    hdr_h = _section_header(slide, x, y, col_w,
                             'Sol 2: Adversarial Disentanglement (Cat B)', '4')
    y += hdr_h

    sol2_h = Emu(4400000)
    _section_body(slide, x, y, col_w, sol2_h)

    sol2_paras = [
        {'text': 'Siamese Encoder (shared weights, 813K params):',
         'font_size': BODY_FONT, 'bold': True, 'font_color': TEAL, 'space_after': 5},
        _bullet_detail('Char embedding (32d)',
                        'Sub-word stylistic patterns (vocab=97)'),
        _bullet_detail('Multi-width CNN (k=3,5,7)',
                        '128 filters each -> 384d, MaxPool(3)'),
        _bullet_detail('BiLSTM (hidden=128)',
                        'Bidirectional -> 256d sequential features'),
        _bullet_detail('Additive Attention',
                        'Bahdanau-style, focuses on salient style cues'),
        _bullet_detail('Projection (128d)',
                        'ReLU + Dropout(0.3) bottleneck'),
        {'text': '', 'font_size': Pt(3), 'space_after': 5},
        {'text': 'Comparison: [v\u2081, v\u2082, |v\u2081\u2212v\u2082|, v\u2081\u2299v\u2082] = 512d',
         'font_size': BODY_FONT, 'bold': True, 'font_color': DARK, 'space_after': 3},
        {'text': 'MLP Classifier: 512 -> 256 -> 64 -> 1 (BCE loss)',
         'font_size': BODY_FONT, 'bold': True, 'font_color': DARK, 'space_after': 6},
        {'text': 'Adversarial Debiasing (Novel):',
         'font_size': BODY_FONT, 'bold': True, 'font_color': CORAL, 'space_after': 5},
        _bullet_detail('Gradient Reversal Layer',
                        'Ganin & Lempitsky (2015)', CORAL, MID_GRAY),
        _bullet_detail('Topic pseudo-labels',
                        'TF-IDF + KMeans (k=10)', CORAL, MID_GRAY),
        _bullet('Lambda ramp: 0 -> 0.05 over 20 epochs; topic_weight=0.02'),
        _bullet('Forces topic-invariant style representations'),
        {'text': '', 'font_size': Pt(3), 'space_after': 5},
        {'text': 'Training: lr=2e-4, batch=64, 40 epochs, char perturbation aug',
         'font_size': SMALL_FONT, 'font_color': MID_GRAY, 'italic': True, 'space_after': 3},
    ]
    _rich_textbox(slide, x + PAD_INNER, y + Emu(110000),
                  col_w - 2 * PAD_INNER, sol2_h - Emu(180000), sol2_paras)
    y += sol2_h + SEC_GAP

    # --- ARCHITECTURE DIAGRAM ---
    hdr_h = _section_header(slide, x, y, col_w, 'Neural Architecture (Sol 2)')
    y += hdr_h

    arch_h = Emu(4800000)
    _section_body(slide, x, y, col_w, arch_h)

    if arch_path.exists():
        img_margin = Emu(100000)
        img_w = col_w - 2 * img_margin
        img_h = Emu(4350000)
        slide.shapes.add_picture(str(arch_path),
                                  x + img_margin, y + Emu(200000),
                                  img_w, img_h)
    y += arch_h + SEC_GAP

    # --- RESULTS TABLE ---
    hdr_h = _section_header(slide, x, y, col_w, 'Results (Dev Set)', '5')
    y += hdr_h

    results_h = Emu(3850000)
    _section_body(slide, x, y, col_w, results_h)

    # Build a clean results table
    results_paras = [
        {'runs': [
            {'text': 'Model', 'font_size': BODY_FONT, 'bold': True, 'font_color': NAVY},
            {'text': '                           ', 'font_size': BODY_FONT},
            {'text': 'F1', 'font_size': BODY_FONT, 'bold': True, 'font_color': NAVY},
            {'text': '          ', 'font_size': BODY_FONT},
            {'text': 'MCC', 'font_size': BODY_FONT, 'bold': True, 'font_color': NAVY},
            {'text': '          ', 'font_size': BODY_FONT},
            {'text': '\u0394 Baseline', 'font_size': BODY_FONT, 'bold': True, 'font_color': NAVY},
        ], 'font_size': BODY_FONT, 'space_after': 5},
        {'text': '\u2500' * 52,
         'font_size': SMALL_FONT, 'font_color': BORDER_COLOR, 'space_after': 5},
        # Baselines
        {'runs': [
            {'text': 'SVM (Cat A baseline)     ', 'font_size': BODY_FONT_SM, 'font_color': MID_GRAY},
            {'text': '0.5610', 'font_size': BODY_FONT_SM, 'font_color': MID_GRAY},
            {'text': '       \u2014             \u2014', 'font_size': BODY_FONT_SM, 'font_color': LIGHT_GRAY}
        ], 'font_size': BODY_FONT_SM, 'space_after': 3},
        {'runs': [
            {'text': 'LSTM (Cat B baseline)   ', 'font_size': BODY_FONT_SM, 'font_color': MID_GRAY},
            {'text': '0.6226', 'font_size': BODY_FONT_SM, 'font_color': MID_GRAY},
            {'text': '       \u2014             \u2014', 'font_size': BODY_FONT_SM, 'font_color': LIGHT_GRAY}
        ], 'font_size': BODY_FONT_SM, 'space_after': 3},
        {'runs': [
            {'text': 'BERT (Cat C baseline)   ', 'font_size': BODY_FONT_SM, 'font_color': MID_GRAY},
            {'text': '0.7854', 'font_size': BODY_FONT_SM, 'font_color': MID_GRAY},
            {'text': '       \u2014             \u2014', 'font_size': BODY_FONT_SM, 'font_color': LIGHT_GRAY}
        ], 'font_size': BODY_FONT_SM, 'space_after': 5},
        {'text': '\u2500' * 52,
         'font_size': SMALL_FONT, 'font_color': BORDER_COLOR, 'space_after': 5},
        # Our solutions
        {'runs': [
            {'text': '\u2605 Sol 1 (Cat A)         ', 'font_size': BODY_FONT, 'bold': True,
             'font_color': ACCENT_GREEN},
            {'text': '0.7340', 'font_size': BODY_FONT, 'bold': True,
             'font_color': ACCENT_GREEN},
            {'text': '   ', 'font_size': BODY_FONT},
            {'text': '0.469', 'font_size': BODY_FONT, 'bold': True,
             'font_color': ACCENT_GREEN},
            {'text': '   ', 'font_size': BODY_FONT},
            {'text': '+30.8% ***', 'font_size': BODY_FONT, 'bold': True,
             'font_color': ACCENT_GREEN}
        ], 'font_size': BODY_FONT, 'space_after': 5},
        {'runs': [
            {'text': '\u2605 Sol 2 (Cat B)         ', 'font_size': BODY_FONT, 'bold': True,
             'font_color': ACCENT_ORANGE},
            {'text': '0.7422', 'font_size': BODY_FONT, 'bold': True,
             'font_color': ACCENT_ORANGE},
            {'text': '   ', 'font_size': BODY_FONT},
            {'text': '0.485', 'font_size': BODY_FONT, 'bold': True,
             'font_color': ACCENT_ORANGE},
            {'text': '   ', 'font_size': BODY_FONT},
            {'text': '+19.2% ***', 'font_size': BODY_FONT, 'bold': True,
             'font_color': ACCENT_ORANGE}
        ], 'font_size': BODY_FONT, 'space_after': 6},
        {'text': '\u2500' * 52,
         'font_size': SMALL_FONT, 'font_color': BORDER_COLOR, 'space_after': 5},
        {'text': '*** p < 0.001, McNemar\u2019s paired test',
         'font_size': SMALL_FONT, 'font_color': MID_GRAY, 'italic': True, 'space_after': 3},
        {'text': 'Both solutions significantly outperform same-category baselines.',
         'font_size': SMALL_FONT, 'font_color': MID_GRAY, 'space_after': 3},
        {'text': 'Neither surpasses BERT (Cat C) \u2014 expected given category constraints.',
         'font_size': SMALL_FONT, 'font_color': MID_GRAY, 'italic': True, 'space_after': 3},
    ]
    _rich_textbox(slide, x + PAD_INNER, y + Emu(110000),
                  col_w - 2 * PAD_INNER, results_h - Emu(180000), results_paras)
    y += results_h + SEC_GAP

    # --- EVALUATION & CONCLUSIONS (moved to Column 2 for balance) ---
    hdr_h = _section_header(slide, x, y, col_w, 'Evaluation & Conclusions', '8')
    y += hdr_h

    eval_h = Emu(2200000)
    _section_body(slide, x, y, col_w, eval_h)

    eval_paras = [
        {'text': 'Evaluation Methods:',
         'font_size': BODY_FONT, 'bold': True, 'font_color': TEAL, 'space_after': 3},
        _bullet('Macro F1 (primary), MCC, per-class P/R'),
        _bullet('McNemar\u2019s paired significance test (p < 0.001)'),
        _bullet('Feature ablation (Sol 1) + GRL ablation (Sol 2)'),
        _bullet('Inter-model agreement (Cohen\u2019s \u03ba \u2248 0.40)'),
        {'text': '', 'font_size': Pt(3), 'space_after': 4},
        {'text': 'Key Conclusions:',
         'font_size': BODY_FONT, 'bold': True, 'font_color': TEAL, 'space_after': 3},
        _bullet('Both solutions significantly outperform their baselines'),
        _bullet('Complementary approaches: feature engineering + neural'),
        _bullet('Novel features and adversarial debiasing add measurable value'),
    ]
    _rich_textbox(slide, x + PAD_INNER, y + Emu(110000),
                  col_w - 2 * PAD_INNER, eval_h - Emu(180000), eval_paras)
    y += eval_h

    # =====================================================
    # COLUMN 3: Charts + Error Analysis + Ethics + References
    # =====================================================
    x = col_x[2]
    y = y_start

    # --- F1 BAR CHART ---
    hdr_h = _section_header(slide, x, y, col_w, 'Performance Comparison', '6')
    y += hdr_h

    f1_h = Emu(3600000)
    _section_body(slide, x, y, col_w, f1_h)

    if f1_path.exists():
        img_margin = Emu(120000)
        img_w = col_w - 2 * img_margin
        img_h = Emu(3150000)
        slide.shapes.add_picture(str(f1_path),
                                  x + img_margin, y + Emu(220000),
                                  img_w, img_h)
    y += f1_h + SEC_GAP

    # --- CONFUSION MATRICES ---
    hdr_h = _section_header(slide, x, y, col_w, 'Confusion Matrices')
    y += hdr_h

    cm_h = Emu(3100000)
    _section_body(slide, x, y, col_w, cm_h)

    if cm_path.exists():
        img_margin = Emu(80000)
        img_w = col_w - 2 * img_margin
        img_h = Emu(2650000)
        slide.shapes.add_picture(str(cm_path),
                                  x + img_margin, y + Emu(220000),
                                  img_w, img_h)
    y += cm_h + SEC_GAP

    # --- ERROR ANALYSIS ---
    hdr_h = _section_header(slide, x, y, col_w, 'Error Analysis', '7')
    y += hdr_h

    error_h = Emu(3200000)
    _section_body(slide, x, y, col_w, error_h)

    error_paras = [
        _kv('Inter-model agreement', '69.8% (Cohen\u2019s \u03ba = 0.396)', TEAL, DARK),
        {'text': '30.2% disagreement \u2014 among those, Sol 1 correct 48.6%, Sol 2 correct 51.4%.',
         'font_size': BODY_FONT_SM, 'font_color': MID_GRAY, 'space_after': 5},
        {'text': 'Per-Class Strengths:',
         'font_size': BODY_FONT, 'bold': True, 'font_color': DARK, 'space_after': 4},
        {'runs': [
            {'text': '  \u2022  Sol 1: ', 'font_size': BODY_FONT_SM, 'bold': True,
             'font_color': ACCENT_GREEN},
            {'text': 'Better at "Different Author" (75.6% vs 73.3%) \u2014 '
                     'stylometric features detect differences well.',
             'font_size': BODY_FONT_SM}
        ], 'font_size': BODY_FONT_SM, 'space_after': 3},
        {'runs': [
            {'text': '  \u2022  Sol 2: ', 'font_size': BODY_FONT_SM, 'bold': True,
             'font_color': ACCENT_ORANGE},
            {'text': 'Better at "Same Author" (75.1% vs 71.3%) \u2014 '
                     'neural model captures stylistic similarity.',
             'font_size': BODY_FONT_SM}
        ], 'font_size': BODY_FONT_SM, 'space_after': 3},
        {'runs': [
            {'text': '  \u2022  Topic confound: ', 'font_size': BODY_FONT_SM, 'bold': True,
             'font_color': CORAL},
            {'text': 'Elevated FP on same-topic pairs; GRL partially mitigates.',
             'font_size': BODY_FONT_SM}
        ], 'font_size': BODY_FONT_SM, 'space_after': 5},
        _kv('Feature ablation', 'Removing 21 novel features: -1.8% F1',
            TEAL, DARK),
        _kv('GRL ablation', 'Removing GRL: -1.2% F1 on topic-confounded pairs',
            TEAL, DARK),
    ]
    _rich_textbox(slide, x + PAD_INNER, y + Emu(110000),
                  col_w - 2 * PAD_INNER, error_h - Emu(180000), error_paras)
    y += error_h + SEC_GAP

    # --- LIMITATIONS & ETHICS ---
    hdr_h = _section_header(slide, x, y, col_w, 'Limitations & Ethics')
    y += hdr_h

    lim_h = Emu(2300000)
    _section_body(slide, x, y, col_w, lim_h)

    lim_paras = [
        {'text': 'Limitations:',
         'font_size': BODY_FONT, 'bold': True, 'font_color': CORAL, 'space_after': 3},
        _bullet('Closed-mode: may not generalize to unseen domains'),
        _bullet('English-only; cultural writing biases not addressed'),
        _bullet('Neither Sol 1 nor Sol 2 surpasses BERT baseline (0.7854)'),
        _bullet('Style verification can be circumvented by imitation attacks'),
        {'text': '', 'font_size': Pt(3), 'space_after': 4},
        {'text': 'Ethical Considerations:',
         'font_size': BODY_FONT, 'bold': True, 'font_color': ACCENT_ORANGE, 'space_after': 3},
        _bullet('Surveillance risk: could be misused for deanonymization'),
        _bullet('False positives in forensic settings cause serious harm'),
        _bullet('Should not be sole evidence in legal proceedings'),
    ]
    _rich_textbox(slide, x + PAD_INNER, y + Emu(110000),
                  col_w - 2 * PAD_INNER, lim_h - Emu(180000), lim_paras)
    y += lim_h + SEC_GAP

    # --- REFERENCES ---
    ref_remaining = bottom_limit - y
    if ref_remaining > Emu(400000):
        hdr_h = _section_header(slide, x, y, col_w, 'Key References')
        y += hdr_h
        ref_h = min(ref_remaining - hdr_h, Emu(1400000))
        _section_body(slide, x, y, col_w, ref_h)
        ref_paras = [
            {'text': '[1] Stamatatos et al. (2023) Diff-Vectors for Authorship Analysis',
             'font_size': SMALL_FONT, 'font_color': MID_GRAY, 'space_after': 2},
            {'text': '[2] Ganin & Lempitsky (2015) Domain Adaptation by Backpropagation',
             'font_size': SMALL_FONT, 'font_color': MID_GRAY, 'space_after': 2},
            {'text': '[3] Abbasi & Chen (2008) Writeprints: Stylometric Identification',
             'font_size': SMALL_FONT, 'font_color': MID_GRAY, 'space_after': 2},
            {'text': '[4] Boenninghoff et al. (2019) AdHominem: Explainable AV',
             'font_size': SMALL_FONT, 'font_color': MID_GRAY, 'space_after': 2},
        ]
        _rich_textbox(slide, x + PAD_INNER, y + Emu(100000),
                      col_w - 2 * PAD_INNER, ref_h - Emu(160000), ref_paras)

    # =====================================================
    # FOOTER
    # =====================================================
    footer_y = SLIDE_H - footer_h
    _shape(slide, 0, footer_y, SLIDE_W, footer_h, NAVY_DARK)
    _shape(slide, 0, footer_y, SLIDE_W, Emu(40000), TEAL)

    _textbox(slide, Emu(500000), footer_y + Emu(100000),
             SLIDE_W - Emu(1000000), footer_h - Emu(160000),
             'COMP34812 Natural Language Understanding  \u2502  '
             'University of Manchester  \u2502  Group 34  \u2502  '
             '2025\u201326  \u2502  Closed Mode  \u2502  '
             'Predictions: Group_34_A.csv + Group_34_B.csv',
             Pt(20), RGBColor(0x99, 0xAA, 0xBB), bold=False,
             alignment=PP_ALIGN.CENTER)

    # =====================================================
    # SAVE
    # =====================================================
    prs.save(str(OUTPUT_PATH))
    print(f"\n{'=' * 60}")
    print(f"  Poster saved to {OUTPUT_PATH}")
    print(f"  Slide dimensions: 33.1\" x 23.4\" (A1 landscape)")
    print(f"{'=' * 60}")


if __name__ == '__main__':
    create_poster()
